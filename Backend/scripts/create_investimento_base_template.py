"""Gera slides/global/investimento_base.pptx.

Cria um template mínimo para o slide de Investimento:
- Título "2. Investimento" (vermelho, bold)
- Shape âncora com texto literal "{{tabela_investimento}}" — o builder
  detecta, lê posição/tamanho, remove e desenha a tabela no lugar.
- Footer "*** Valores válidos somente para contratação simultânea dos serviços."

Reexecutar este script sobrescreve o arquivo. Idempotente.
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt

REPO_ROOT = Path(__file__).resolve().parents[2]
OUTPUT = REPO_ROOT / "Backend/services/pptx-generator-service/slides/global/investimento_base.pptx"

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

VERMELHO = RGBColor(0xD4, 0x00, 0x00)
PRETO    = RGBColor(0x1C, 0x1C, 0x1C)
CINZA    = RGBColor(0x66, 0x66, 0x66)


def _add_text(slide, text, *, x, y, w, h, size, color, bold=False, anchor=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return box


def main() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Título
    _add_text(
        slide, "2. Investimento",
        x=0.5, y=0.25, w=12.0, h=0.8,
        size=32, color=VERMELHO, bold=True,
    )

    # Âncora da tabela (será removida e substituída pela tabela real)
    _add_text(
        slide, "{{tabela_investimento}}",
        x=0.5, y=1.2, w=12.33, h=5.5,
        size=10, color=RGBColor(0xCC, 0xCC, 0xCC),
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # Footer
    _add_text(
        slide, "*** Valores válidos somente para contratação simultânea dos serviços.",
        x=0.5, y=6.95, w=12.33, h=0.35,
        size=10, color=CINZA, bold=False,
    )

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"OK: {OUTPUT}")


if __name__ == "__main__":
    main()
