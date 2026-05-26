import logging
import os
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from placeholder_engine import _replace_placeholders
from slide_copier import _copy_slide

from .catalog import InvestItem, TextRun, get_items

_PLACEHOLDER_RE = re.compile(r"\{\{\s*(\w+)\s*\}\}")

logger = logging.getLogger("pptx_generator.investimento")

ANCHOR_TOKEN = "{{tabela_investimento}}"

_HEADER_BG = RGBColor(0x80, 0x80, 0x80)
_HEADER_FG = RGBColor(0xFF, 0xFF, 0xFF)
_ZEBRA_BG  = RGBColor(0xF2, 0xF2, 0xF2)
_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
_BODY_FG   = RGBColor(0x00, 0x00, 0x00)
_VERMELHO  = RGBColor(0xFF, 0x00, 0x00)

_COLUMN_WIDTH_PCT = (0.55, 0.13, 0.13, 0.19)
_HEADER_LABELS    = ("Descrição", "Qtde.", "Unid.", "Valor")
_HEADER_FONT_PT   = 14
_BODY_FONT_PT     = 12
_FONT_NAME        = "Poppins"
_TOTAL_LABEL      = "VALOR TOTAL"
_TOTAL_VALUE      = "R$"

# Altura máxima do cabeçalho/total. Sem isso, python-pptx divide a altura total
# da tabela igualmente entre as linhas — com poucos itens o header fica gigante.
_HEADER_ROW_HEIGHT = Inches(0.5)
_TOTAL_ROW_HEIGHT  = Inches(0.5)

# Bordas: contorno externo + horizontais entre rows (sem verticais internas).
_BORDER_COLOR_HEX = "000000"
_BORDER_WIDTH_EMU = int(Pt(1))  # 12700 EMU = 1pt


def compose_investimento(
    merged: Presentation,
    base_path: str,
    product_id: str,
    variant_id: str,
    values: dict,
    ctx: dict,
    img_counter: list[int],
) -> None:
    if not os.path.exists(base_path):
        logger.warning("investimento_base não encontrado: %s", base_path)
        return

    items = get_items(product_id, variant_id, values)

    src_base = Presentation(base_path)
    new_slide = _copy_slide(merged, src_base.slides[0], img_counter)

    anchor = _find_and_remove_anchor(new_slide)
    if anchor is None:
        logger.warning("âncora '%s' não encontrada em %s", ANCHOR_TOKEN, base_path)
        _replace_placeholders(new_slide, ctx)
        return

    _replace_placeholders(new_slide, ctx)

    if not items:
        logger.info("nenhum item de investimento aplicável", extra={"product_id": product_id})
        return

    left, top, width, height = anchor
    n_rows = len(items) + 2  # header + items + total
    table_shape = new_slide.shapes.add_table(n_rows, 4, left, top, width, height)
    table = table_shape.table

    _apply_column_widths(table, width)
    _apply_row_heights(table, height, n_rows)
    _build_header_row(table)
    for i, item in enumerate(items, start=1):
        _build_item_row(table, i, item, values, ctx)
    _build_total_row(table, n_rows - 1)
    _apply_table_borders(table, n_rows)


# ── internals ────────────────────────────────────────────────────────────────

def _find_and_remove_anchor(slide):
    sp_tree = slide.shapes._spTree
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.text_frame.text.strip() != ANCHOR_TOKEN:
            continue
        coords = (shape.left, shape.top, shape.width, shape.height)
        sp_tree.remove(shape.element)
        return coords
    return None


def _apply_row_heights(table, total_height: int, n_rows: int) -> None:
    """Fixa altura do header e do total, distribui o restante entre os itens.

    Sem isso, python-pptx divide `total_height` igualmente entre as N linhas:
    com 3 linhas o header fica enorme; com 10, fica esticado. O usuário/PowerPoint
    pode auto-crescer células se o conteúdo exigir, mas a altura definida é o piso.
    """
    body_rows = max(n_rows - 2, 0)  # header e total fora
    remaining = max(total_height - _HEADER_ROW_HEIGHT - _TOTAL_ROW_HEIGHT, 0)
    body_height = (remaining // body_rows) if body_rows else 0

    table.rows[0].height = _HEADER_ROW_HEIGHT
    for r_idx in range(1, n_rows - 1):
        table.rows[r_idx].height = body_height
    table.rows[n_rows - 1].height = _TOTAL_ROW_HEIGHT


def _apply_column_widths(table, total_width: int) -> None:
    widths = [int(total_width * pct) for pct in _COLUMN_WIDTH_PCT]
    drift = total_width - sum(widths)
    widths[0] += drift  # corrige arredondamento na coluna maior
    for col, w in zip(table.columns, widths):
        col.width = w


_BORDER_EDGE_ORDER = ("L", "R", "T", "B")  # ordem exigida pelo schema OOXML em <a:tcPr>


def _make_border_element(edge: str):
    xml = (
        f'<a:ln{edge} xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        f'w="{_BORDER_WIDTH_EMU}" cap="flat" cmpd="sng" algn="ctr">'
        f'<a:solidFill><a:srgbClr val="{_BORDER_COLOR_HEX}"/></a:solidFill>'
        f'<a:prstDash val="solid"/>'
        f'</a:ln{edge}>'
    )
    return parse_xml(xml)


def _set_cell_border(cell, edges: set[str]) -> None:
    """Aplica bordas nas arestas indicadas (subconjunto de {L,R,T,B}).

    Os elementos lnL/lnR/lnT/lnB devem vir no início do <a:tcPr>, antes do
    <a:solidFill> de preenchimento. Inserindo na ordem reversa no índice 0,
    o resultado final fica L,R,T,B (ordem do schema).
    """
    tcPr = cell._tc.get_or_add_tcPr()

    # Idempotência: remove bordas pré-existentes
    for edge in _BORDER_EDGE_ORDER:
        for el in tcPr.findall(qn(f"a:ln{edge}")):
            tcPr.remove(el)

    for edge in reversed(_BORDER_EDGE_ORDER):  # B, T, R, L
        if edge in edges:
            tcPr.insert(0, _make_border_element(edge))


def _apply_table_borders(table, n_rows: int, n_cols: int = 4) -> None:
    """Contorno externo + horizontais entre rows; sem verticais internas."""
    for r in range(n_rows):
        for c in range(n_cols):
            edges = {"T", "B"}
            if c == 0:
                edges.add("L")
            if c == n_cols - 1:
                edges.add("R")
            _set_cell_border(table.cell(r, c), edges)


def _build_header_row(table) -> None:
    for col_idx, label in enumerate(_HEADER_LABELS):
        cell = table.cell(0, col_idx)
        _set_cell_fill(cell, _HEADER_BG)
        align = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.CENTER
        _set_cell_text(
            cell,
            [TextRun(text=label, bold=True, color=_HEADER_FG, size_pt=_HEADER_FONT_PT)],
            align=align,
        )


def _build_item_row(table, row_idx: int, item: InvestItem, values: dict, ctx: dict) -> None:
    bg = _ZEBRA_BG if row_idx % 2 == 1 else _WHITE
    for col in range(4):
        _set_cell_fill(table.cell(row_idx, col), bg)

    # Coluna 0 — descrição (rich text, placeholders {{ key }} resolvidos via ctx)
    runs = item.resolve_runs(values)
    runs = [_with_default_size(r, _BODY_FONT_PT) for r in runs]
    runs = _resolve_runs(runs, ctx)
    _set_cell_text(table.cell(row_idx, 0), runs, align=PP_ALIGN.LEFT)

    # Coluna 1 — quantidade
    try:
        qtde = float(item.qtde_resolver(values) or 0)
    except (TypeError, ValueError):
        qtde = 0.0
    qtde_fmt = f"{qtde:.2f}".replace(".", ",")
    _set_cell_text(
        table.cell(row_idx, 1),
        [TextRun(text=qtde_fmt, color=_BODY_FG, size_pt=_BODY_FONT_PT)],
        align=PP_ALIGN.CENTER,
    )

    # Coluna 2 — unidade
    _set_cell_text(
        table.cell(row_idx, 2),
        [TextRun(text=item.unidade, color=_BODY_FG, size_pt=_BODY_FONT_PT)],
        align=PP_ALIGN.CENTER,
    )

    # Coluna 3 — valor (deixado em branco para preenchimento manual)
    _set_cell_text(
        table.cell(row_idx, 3),
        [TextRun(text="R$", color=_BODY_FG, size_pt=_BODY_FONT_PT)],
        align=PP_ALIGN.CENTER,
    )


def _build_total_row(table, row_idx: int) -> None:
    for col in range(4):
        _set_cell_fill(table.cell(row_idx, col), _WHITE)

    _set_cell_text(
        table.cell(row_idx, 0),
        [TextRun(text=_TOTAL_LABEL, bold=True, color=_VERMELHO, size_pt=_HEADER_FONT_PT)],
        align=PP_ALIGN.LEFT,
    )
    for col in (1, 2):
        _set_cell_text(table.cell(row_idx, col), [TextRun(text="", size_pt=_HEADER_FONT_PT)])
    _set_cell_text(
        table.cell(row_idx, 3),
        [TextRun(text=_TOTAL_VALUE, bold=True, color=_VERMELHO, size_pt=_HEADER_FONT_PT)],
        align=PP_ALIGN.CENTER,
    )


def _set_cell_fill(cell, rgb: RGBColor) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = rgb


def _set_cell_text(cell, runs: list[TextRun], *, align: PP_ALIGN = PP_ALIGN.LEFT) -> None:
    tf = cell.text_frame
    tf.word_wrap = True
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Limpa qualquer parágrafo/run existente reaproveitando o primeiro parágrafo.
    first_para = tf.paragraphs[0]
    _clear_paragraph(first_para)
    # remove parágrafos extras (caso o python-pptx tenha criado mais de um)
    for extra in list(tf.paragraphs[1:]):
        extra._p.getparent().remove(extra._p)

    target_para = first_para
    target_para.alignment = align

    for run_spec in runs:
        for i, chunk in enumerate(run_spec.text.split("\n")):
            if i > 0:
                target_para = tf.add_paragraph()
                target_para.alignment = align
            if chunk == "":
                continue
            r = target_para.add_run()
            r.text = chunk
            r.font.name = _FONT_NAME
            if run_spec.bold:
                r.font.bold = True
            if run_spec.italic:
                r.font.italic = True
            if run_spec.color is not None:
                r.font.color.rgb = run_spec.color
            if run_spec.size_pt is not None:
                r.font.size = Pt(run_spec.size_pt)


def _clear_paragraph(paragraph) -> None:
    for run in list(paragraph.runs):
        run._r.getparent().remove(run._r)


def _resolve_placeholders(text: str, ctx: dict) -> str:
    """Substitui `{{ key }}` / `{{key}}` por ctx[key]. Chaves ausentes viram vazio."""
    if "{{" not in text:
        return text

    def repl(m: re.Match) -> str:
        v = ctx.get(m.group(1))
        return str(v) if v is not None else ""

    return _PLACEHOLDER_RE.sub(repl, text)


def _resolve_runs(runs: list[TextRun], ctx: dict) -> list[TextRun]:
    return [
        TextRun(
            text=_resolve_placeholders(r.text, ctx),
            bold=r.bold,
            italic=r.italic,
            color=r.color,
            size_pt=r.size_pt,
        )
        for r in runs
    ]


def _with_default_size(run: TextRun, size_pt: int) -> TextRun:
    if run.size_pt is not None:
        return run
    return TextRun(
        text=run.text,
        bold=run.bold,
        italic=run.italic,
        color=run.color if run.color is not None else _BODY_FG,
        size_pt=size_pt,
    )
