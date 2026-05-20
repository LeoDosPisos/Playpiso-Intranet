import copy
import logging
import os
import time
from datetime import datetime
from io import BytesIO

logger = logging.getLogger("pptx_generator.merger")

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.parts.image import ImagePart
from pptx.opc.packuri import PackURI
from pptx.util import Inches, Pt

_R_EMBED = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'
_img_counter = 0

_PRETO = RGBColor(0x1C, 0x1C, 0x1C)
# Placeholders whose replacement text must be rendered in black
_RESET_COLOR_KEYS = frozenset({"sumario"})

_GALVANIZACAO_LABELS = {'fogo': 'a fogo', 'eletrolitico': 'eletroliticamente'}

_SISTEMA_ALAMBRADO_LABELS = {'gaiola': 'Gaiola', 'trapezio': 'Trapézio'}

_ILUMINACAO_FALLBACK_KEYS = (
    'quantidade_projetores', 'potencia_projetores',
    'quantidade_postes_iluminacao', 'altura_postes_iluminacao',
    'quantidade_cruzetas',
)

_ESTRUTURA_BASQUETE_LABELS = {'metalica': 'Metálica', 'hidraulica': 'Hidráulica', 'comum': 'Comum'}

# Ordem canônica para exibição no slide
_TRAVAMENTO_ORDER = ['travamento_superior', 'travamento_intermediario', 'travamento_inferior']
_TRAVAMENTO_LABELS = {
    'travamento_superior':      'superior',
    'travamento_intermediario': 'intermediário',
    'travamento_inferior':      'inferior',
}

def _is_truthy(value) -> bool:
    """Normaliza valores que podem vir como bool, str ou None do contexto.

    Necessário porque _build_context converte tudo com str(v), então
    False vira "False" (string truthy).
    """
    if isinstance(value, bool):
        return value
    if isinstance(value, str):
        return value.lower() in ('true', '1', 'yes')
    return bool(value)


SLIDES_DIR = os.path.join(os.path.dirname(__file__), "slides")
SLIDE_W    = Inches(13.33)
SLIDE_H    = Inches(7.5)

# Maps slideId (from slideRegistry.ts) → path relative to SLIDES_DIR
SLIDE_FILE_MAP: dict[str, str] = {
    # ── Global pré-produto ────────────────────────────────────────────────
    "capa":                               "global/capa.pptx",
    "portfolio":                          "global/portfolio.pptx",
    "sobre_empresa":                      "global/sobre_empresa.pptx",
    "pilares":                            "global/pilares.pptx",
    "parceiros":                          "global/parceiros.pptx",
    "dados_cliente":                      "global/dados_cliente.pptx",
    "sumario":                            "global/sumario.pptx",
    # ── beach_tenis ───────────────────────────────────────────────────────
    "hero_beach_tenis":                   "beach_tenis/hero.pptx",
    "areia_rio_beach_tenis":              "beach_tenis/areia_rio.pptx",
    "areia_quartzo_beach_tenis":          "beach_tenis/areia_quartzo.pptx",
    "protecao_eva_beach_tenis":           "beach_tenis/protecao_eva.pptx",
    # fechamentos_beach_tenis: resolvido dinamicamente via compose_fechamentos
    "acessorio_beach_tenis":              "beach_tenis/acessorio.pptx",
    "investimento_beach_tenis":           "beach_tenis/investimento.pptx",
    # ── quadra_tenis ──────────────────────────────────────────────────────
    "hero_piso_asfaltico_quadra_tenis":   "quadra_tenis/hero_piso_asfaltico.pptx",
    "hero_saibro_quadra_tenis":           "quadra_tenis/hero_saibro.pptx",
    "hero_grama_quadra_tenis":            "quadra_tenis/hero_grama.pptx",
    "specs_piso_asfaltico":               "quadra_tenis/specs_piso_asfaltico.pptx",
    "specs_saibro":                       "quadra_tenis/specs_saibro.pptx",
    "specs_grama":                        "quadra_tenis/specs_grama.pptx",
    "playcushion_quadra_tenis":           "quadra_tenis/playcushion.pptx",
    # fechamentos_quadra_tenis: resolvido dinamicamente via compose_fechamentos
    "cores_piso_asfaltico":               "quadra_tenis/cores_piso_asfaltico.pptx",
    "detalhe_construtivo_sem_playcushion":"quadra_tenis/detalhe_construtivo_sem_playcushion.pptx",
    "detalhe_construtivo_com_playcushion":"quadra_tenis/detalhe_construtivo.pptx",
    "investimento_piso_asfaltico_quadra_tenis": "quadra_tenis/investimento_piso_asfaltico.pptx",
    "investimento_saibro_quadra_tenis":         "quadra_tenis/investimento_saibro.pptx",
    "investimento_grama_quadra_tenis":          "quadra_tenis/investimento_grama.pptx",
    # ── quadra_poliesportiva ──────────────────────────────────────────────
    "hero_piso_asfaltico_quadra_poliesportiva":                 "quadra_poli/hero_piso_asfaltico.pptx",
    "specs_piso_asfaltico_quadra_poliesportiva":                "quadra_poli/specs_piso_asfaltico.pptx",
    # acessorios_quadra_poliesportiva: resolvido dinamicamente via compose_acessorios
    # fechamentos_quadra_poliesportiva: resolvido dinamicamente via compose_fechamentos
    "cores_piso_asfaltico_quadra_poliesportiva":                "quadra_poli/cores_piso_asfaltico.pptx",
    "investimento_piso_asfaltico_quadra_poliesportiva":         "quadra_poli/investimento_piso_asfaltico.pptx",
    # ── Global pós-produto ────────────────────────────────────────────────
    "condicoes_pagamento_direto_a":       "global/condicoes_pagamento_direto_a.pptx",
    "condicoes_pagamento_direto_b":       "global/condicoes_pagamento_direto_b.pptx",
    "condicoes_pagamento_playpiso":       "global/condicoes_pagamento_playpiso.pptx",
    "prazos_garantia":                    "global/prazos_garantia.pptx",
    "regras_contratada":                  "global/regras_contratada.pptx",
    "regras_contratante":                 "global/regras_contratante.pptx",
    "consideracoes_gerais":               "global/consideracoes_gerais.pptx",
    "encerramento":                       "global/encerramento.pptx",
}


# Mapeia slideId → [(título_da_seção, chave_no_contexto), ...]
# Seções cujo valor de contexto seja falsy são removidas antes da substituição.
CONDITIONAL_SECTIONS: dict[str, list[tuple[str, str]]] = {}


def _slide_template_path(slide_id: str) -> str | None:
    rel = SLIDE_FILE_MAP.get(slide_id)
    if not rel:
        return None
    path = os.path.join(SLIDES_DIR, rel)
    return path if os.path.exists(path) else None


def is_slide_available(slide_id: str) -> bool:
    return _slide_template_path(slide_id) is not None


def get_available_slides() -> list[str]:
    return [sid for sid in SLIDE_FILE_MAP if is_slide_available(sid)]


_PRODUCT_SLIDES_DIR: dict[str, str] = {
    "beach_tenis":          "beach_tenis",
    "quadra_tenis":         "quadra_tenis",
    "quadra_poliesportiva": "quadra_poli",
}

_FECHAMENTOS_SECTIONS: list[tuple[str, str]] = [
    ("possui_alambrado",         "alambrado"),
    ("possui_iluminacao",        "iluminacao"),
    ("possui_tela_superior",     "tela_superior"),
    ("possui_tela_sombreamento", "tela_sombreamento"),
]

_CM_TO_EMU                    = 914400 / 2.54
_CONTENT_TOP_EMU              = round(0.5 * _CM_TO_EMU)    # fechamentos
_ACESSORIOS_CONTENT_TOP_EMU   = round(2.25 * _CM_TO_EMU)   # abaixo do título "Acessórios" (bot=2.52cm)
_SECTION_GAP_EMU              = round(-0.5 * _CM_TO_EMU)   # fechamentos
_ACESSORIOS_SECTION_GAP_EMU   = round(-0.1 * _CM_TO_EMU)    # espaço entre sub-seções de acessórios


def _shift_shape_top(el, delta_emu: int) -> None:
    from pptx.oxml.ns import qn
    for tag in (qn("p:spPr"), qn("p:grpSpPr")):
        spPr = el.find(tag)
        if spPr is None:
            continue
        xfrm = spPr.find(qn("a:xfrm"))
        if xfrm is None:
            continue
        off = xfrm.find(qn("a:off"))
        if off is not None:
            off.set("y", str(int(off.get("y", "0")) + delta_emu))
        return


def compose_fechamentos(
    merged: Presentation,
    base_path: str,
    product_id: str,
    values: dict,
    ctx: dict,
) -> None:
    active = [
        name for field, name in _FECHAMENTOS_SECTIONS if _is_truthy(values.get(field))
    ]
    if not active:
        return

    pages = [active] if len(active) <= 3 else [active[:2], active[2:]]
    product_subdir = _PRODUCT_SLIDES_DIR.get(product_id, product_id)

    for page_sections in pages:
        if not os.path.exists(base_path):
            logger.warning("fechamentos_base não encontrado: %s", base_path)
            return

        src_base = Presentation(base_path)
        new_slide = _copy_slide(merged, src_base.slides[0])

        sp_tree = new_slide.shapes._spTree
        current_top = _CONTENT_TOP_EMU
        # 2 seções por página (caso 4 seções → 2 slides): sem sobreposição
        gap = 0 if len(page_sections) == 2 else _SECTION_GAP_EMU

        for section_name in page_sections:
            comp_path = os.path.join(SLIDES_DIR, product_subdir, f"secao_{section_name}.pptx")
            if not os.path.exists(comp_path):
                logger.warning("componente não encontrado: %s", comp_path)
                continue
            src_comp = Presentation(comp_path)
            comp_slide = src_comp.slides[0]

            # mede a altura real do componente (borda inferior mais baixa)
            comp_height = max(
                (shape.top + shape.height for shape in comp_slide.shapes),
                default=0,
            )

            for shape in comp_slide.shapes:
                el = copy.deepcopy(shape._element)
                _shift_shape_top(el, current_top)
                sp_tree.append(el)

            current_top += comp_height + gap

        _replace_placeholders(new_slide, ctx)


def _get_active_acessorios_sections(values: dict) -> list[str]:
    sections: list[str] = []

    if _is_truthy(values.get("possui_basquete_adulto")):
        estrutura = str(values.get("estrutura_basquete_adulto") or "comum")
        if estrutura in ("metalica", "hidraulica", "comum"):
            sections.append(f"basquete_adulto_{estrutura}")
        else:
            sections.append("basquete_adulto_comum")

    if _is_truthy(values.get("possui_basquete_juvenil")):
        sections.append("basquete_juvenil")

    if _is_truthy(values.get("possui_volei")):
        sections.append("volei")

    if _is_truthy(values.get("possui_tenis")):
        sections.append("tenis")

    if _is_truthy(values.get("possui_futebol_futsal")):
        tipo = str(values.get("tipo_futsal") or "padrao")
        sections.append(f"futsal_{tipo}" if tipo in ("padrao", "mini_trave") else "futsal_padrao")

    return sections


def compose_acessorios(
    merged: Presentation,
    base_path: str,
    product_id: str,
    values: dict,
    ctx: dict,
) -> None:
    active = _get_active_acessorios_sections(values)
    if not active:
        return

    product_subdir = _PRODUCT_SLIDES_DIR.get(product_id, product_id)
    pages = [active[i:i + 4] for i in range(0, len(active), 4)]

    for page_sections in pages:
        if not os.path.exists(base_path):
            logger.warning("acessorios_base não encontrado: %s", base_path)
            return

        src_base = Presentation(base_path)
        new_slide = _copy_slide(merged, src_base.slides[0])
        sp_tree = new_slide.shapes._spTree
        current_top = _ACESSORIOS_CONTENT_TOP_EMU

        for section_name in page_sections:
            comp_path = os.path.join(
                SLIDES_DIR, product_subdir, f"secao_acessorio_{section_name}.pptx"
            )
            if not os.path.exists(comp_path):
                logger.warning("componente acessório não encontrado: %s", comp_path)
                continue

            src_comp = Presentation(comp_path)
            comp_slide = src_comp.slides[0]
            comp_height = max(
                (shape.top + shape.height for shape in comp_slide.shapes),
                default=0,
            )
            for shape in comp_slide.shapes:
                el = copy.deepcopy(shape._element)
                _shift_shape_top(el, current_top)
                sp_tree.append(el)
            current_top += comp_height + _ACESSORIOS_SECTION_GAP_EMU

        _replace_placeholders(new_slide, ctx)


def build_presentation(req) -> bytes:
    t0 = time.perf_counter()

    use_rich_slides = hasattr(req, 'slides') and req.slides is not None
    context = _build_context(req.globalValues, req.productGroups)
    groups_by_index = {i: g for i, g in enumerate(req.productGroups)}

    merged = Presentation()
    merged.slide_width = SLIDE_W
    merged.slide_height = SLIDE_H

    if use_rich_slides:
        for slide_entry in req.slides:
            slide_id = slide_entry.slideId
            dynamic   = slide_entry.dynamic
            group_idx = slide_entry.groupIndex

            if dynamic == "fechamentos":
                product_id = slide_id.removeprefix("fechamentos_")
                group  = groups_by_index.get(group_idx or 0)
                values = dict(group.values) if group else {}
                base_rel  = slide_entry.templateFile.removeprefix("slides/")
                base_path = os.path.join(SLIDES_DIR, base_rel)
                compose_fechamentos(merged, base_path, product_id, values, context)
            elif dynamic == "acessorios":
                product_id = slide_id.removeprefix("acessorios_")
                group  = groups_by_index.get(group_idx or 0)
                values = dict(group.values) if group else {}
                base_rel  = slide_entry.templateFile.removeprefix("slides/")
                base_path = os.path.join(SLIDES_DIR, base_rel)
                compose_acessorios(merged, base_path, product_id, values, context)
            else:
                base_rel  = slide_entry.templateFile.removeprefix("slides/")
                file_path = os.path.join(SLIDES_DIR, base_rel)
                if os.path.exists(file_path):
                    _add_from_file_with_replacement(merged, file_path, slide_id, context)
                else:
                    _add_placeholder_slide(merged, slide_id)
    else:
        for slide_id in (req.slideIds or []):
            file_path = _slide_template_path(slide_id)
            if file_path:
                _add_from_file_with_replacement(merged, file_path, slide_id, context)
            else:
                _add_placeholder_slide(merged, slide_id)

    output = BytesIO()
    merged.save(output)
    result = output.getvalue()

    logger.info(
        "presentation_built",
        extra={
            "slide_count": len(merged.slides),
            "size_bytes": len(result),
            "duration_ms": round((time.perf_counter() - t0) * 1000, 1),
        },
    )
    return result


_NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'


def _copy_background(src_slide, dest_slide) -> None:
    src_cSld = src_slide._element.find(f'{{{_NS_P}}}cSld')
    src_bg = src_cSld.find(f'{{{_NS_P}}}bg') if src_cSld is not None else None

    if src_bg is None:
        return

    dest_cSld = dest_slide._element.find(f'{{{_NS_P}}}cSld')
    existing_bg = dest_cSld.find(f'{{{_NS_P}}}bg')
    if existing_bg is not None:
        dest_cSld.remove(existing_bg)

    dest_cSld.insert(0, copy.deepcopy(src_bg))


def _copy_slide(merged: Presentation, src_slide) -> any:
    global _img_counter
    blank_layout = merged.slide_layouts[6]
    new_slide = merged.slides.add_slide(blank_layout)

    rId_map = {}
    for rId, rel in src_slide.part.rels.items():
        if '/image' in rel.reltype:
            _img_counter += 1
            src_img = rel.target_part
            ext = os.path.splitext(str(src_img.partname))[1]
            new_partname = PackURI(f'/ppt/media/m{_img_counter}{ext}')
            new_img = ImagePart(new_partname, src_img.content_type, merged.part.package, src_img.blob)
            new_rId = new_slide.part.relate_to(new_img, rel.reltype)
            rId_map[rId] = new_rId

    spTree_copy = copy.deepcopy(src_slide.shapes._spTree)
    for el in spTree_copy.iter():
        old = el.get(_R_EMBED)
        if old and old in rId_map:
            el.set(_R_EMBED, rId_map[old])

    sp = new_slide.shapes._spTree
    sp.clear()
    for el in spTree_copy:
        sp.append(el)

    _copy_background(src_slide, new_slide)

    return new_slide


def _remove_section(slide, section_title: str) -> bool:
    """Remove título + conteúdo de uma seção identificada pelo texto exato do título.

    Localiza o título pelo texto, depois remove o shape de conteúdo mais próximo
    abaixo dele (menor top > title_top). Não depende de nomes ou IDs do PowerPoint.
    """
    spTree = slide.shapes._spTree

    title_el = None
    title_top = None
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip() == section_title:
            title_el = shape.element
            title_top = shape.top
            break

    if title_el is None:
        logger.warning("section_not_found", extra={"section": section_title})
        return False

    content_el = None
    min_dist = float('inf')
    for shape in slide.shapes:
        if shape.element is title_el:
            continue
        if shape.has_text_frame and shape.top > title_top:
            dist = shape.top - title_top
            if dist < min_dist:
                min_dist = dist
                content_el = shape.element

    spTree.remove(title_el)
    if content_el is not None:
        spTree.remove(content_el)

    logger.info("section_removed", extra={"section": section_title})
    return True


def _add_from_file_with_replacement(
    merged: Presentation, file_path: str, label: str, context: dict
) -> None:
    src = Presentation(file_path)
    for src_slide in src.slides:
        new_slide = _copy_slide(merged, src_slide)

        if label in CONDITIONAL_SECTIONS:
            for section_title, context_key in CONDITIONAL_SECTIONS[label]:
                if not _is_truthy(context.get(context_key)):
                    _remove_section(new_slide, section_title)

        _replace_placeholders(new_slide, context)

    logger.info("slide_loaded", extra={"label": label, "slide_count": len(src.slides)})


def _fmt_date(iso: str) -> str:
    try:
        return datetime.strptime(iso, "%Y-%m-%d").strftime("%d/%m/%Y")
    except (ValueError, TypeError):
        return iso or ""


def _fmt_dimension(value) -> str:
    try:
        return f"{float(value):.2f}".replace('.', ',') + 'm'
    except (ValueError, TypeError):
        return "—"


def _fmt_numero(value) -> str:
    try:
        return f"{float(value):.2f}".replace('.', ',')
    except (ValueError, TypeError):
        return "—"


def _fmt_alambrado_descricao(values: dict) -> str:
    sistema = values.get('sistema_alambrado', '')
    h_fun = values.get('altura_alambrado_fundos')
    h_lat = values.get('altura_alambrado_laterais')
    if sistema == 'trapezio' and h_fun is not None and h_lat is not None:
        return (
            f"Sistema trapézio: alambrado com fundo de {_fmt_dimension(h_fun)}m"
            f" e corrimão (altura de 1,00m) conectando as laterais de {_fmt_dimension(h_lat)}m;"
        )
    if sistema == 'gaiola' and h_fun is not None:
        return f"Sistema gaiola: alambrado com fundo e laterais de {_fmt_dimension(h_fun)}m;"
    return '—'


def _fmt_travamento(value) -> str:
    if isinstance(value, list):
        items = value
    elif isinstance(value, str) and value:
        items = [t.strip() for t in value.split(',')]
    else:
        return 'Sem travamento'

    selected = [t for t in _TRAVAMENTO_ORDER if t in items]
    if not selected or selected == ['sem_travamento'] or items == ['sem_travamento']:
        return 'Sem travamento'

    labels = [_TRAVAMENTO_LABELS[t] for t in selected]
    if len(labels) == 1:
        return f'Instalação de travamento {labels[0]}'
    return f'Instalação de travamento {", ".join(labels[:-1])} e {labels[-1]}'


def _build_context(global_values: dict, product_groups: list) -> dict:
    ctx = {
        "nome_razao_social": global_values.get("nome_razao_social", ""),
        "nome_contato":      global_values.get("nome_contato", ""),
        "endereco_obra":     global_values.get("endereco_obra", ""),
        "local_obra":        global_values.get("local_obra", ""),
        "telefone":          global_values.get("telefone", ""),
        "email":             global_values.get("email", ""),
        "numero_proposta":   global_values.get("numero_proposta", ""),
        "data_solicitacao":  _fmt_date(global_values.get("data_solicitacao", "")),
        "data_envio":        _fmt_date(global_values.get("data_envio", "")),
        "sumario":           (
            product_groups[0].sumarioText
            if len(product_groups) == 1
            else "\n".join(f"{i + 1}. {g.sumarioText}" for i, g in enumerate(product_groups))
        ),
    }
    for group in product_groups:
        ctx.update({k: str(v) if v is not None else "" for k, v in group.values.items()})
        values = group.values or {}
        if ctx.get('potencia_projetores') == 'outro':
            ctx['potencia_projetores'] = ctx.get('especificar_potencia_projetores') or '—'
        galv = ctx.get('galvanizacao', '')
        if galv in _GALVANIZACAO_LABELS:
            ctx['galvanizacao'] = _GALVANIZACAO_LABELS[galv]
        elif galv == 'outro':
            ctx['galvanizacao'] = ctx.get('especificar_galvanizacao') or '—'
        altura = values.get('altura_portoes')
        largura = values.get('largura_portoes')
        qtd = values.get('quantidade_portoes', 0)
        if qtd and altura is not None and largura is not None:
            ctx['dimensoes_portoes'] = f"{_fmt_dimension(altura)} x {_fmt_dimension(largura)}"
        else:
            ctx['dimensoes_portoes'] = "—"
        ctx['travamento_descricao'] = _fmt_travamento(values.get('travamento'))
        ctx['alambrado_descricao'] = _fmt_alambrado_descricao(values)
        ctx['descricao_alambrado'] = ctx['alambrado_descricao']

        # ── slide investimento ────────────────────────────────────────────
        ctx['quantity']       = _fmt_numero(group.quantity)
        ctx['area_total_fmt'] = _fmt_numero(values.get('area_total'))

        if values.get('possui_alambrado'):
            c_lat = values.get('comprimento_alambrado_laterais')
            h_lat = values.get('altura_alambrado_laterais')
            c_fun = values.get('comprimento_alambrado_fundos')
            h_fun = values.get('altura_alambrado_fundos')
            try:
                area = float(c_lat or 0) * float(h_lat or 0) + float(c_fun or 0) * float(h_fun or 0)
                ctx['area_alambrado'] = _fmt_numero(area) if area else '—'
            except (TypeError, ValueError):
                ctx['area_alambrado'] = '—'
        else:
            ctx['area_alambrado'] = '—'

        ctx['qtde_iluminacao'] = '1,00' if values.get('possui_iluminacao') else '—'

        ctx['area_playcushion'] = (
            _fmt_numero(values.get('area_total')) if values.get('possui_playcushion') else '—'
        )

        # aliases para templates que usam nomes alternativos
        ctx['area_fmt']          = ctx.get('area_total_fmt', '—')
        ctx['area_tela_superior'] = ctx.get('area_total_fmt', '—')

        # kit saibro: "1,00" quando incluso, "—" quando não
        ctx['kit_saibro'] = '1,00' if values.get('possui_kit_saibro') else '—'

        # ── sistema alambrado ─────────────────────────────────────────────
        raw_sistema = values.get('sistema_alambrado', '')
        if raw_sistema:
            ctx['sistema_alambrado'] = _SISTEMA_ALAMBRADO_LABELS.get(raw_sistema, raw_sistema)
        else:
            ctx['sistema_alambrado'] = '—'
        # alias para typo no template investimento_piso_asfaltico (quadra_poli)
        ctx['sistema_alabrado'] = ctx['sistema_alambrado']

        # ── alambrado: fallback quando não selecionado ────────────────────
        if not ctx.get('galvanizacao'):
            ctx['galvanizacao'] = '—'

        # ── iluminação: fallback quando não selecionada ───────────────────
        for _k in _ILUMINACAO_FALLBACK_KEYS:
            if not ctx.get(_k):
                ctx[_k] = '—'

        # ── acessórios esportivos (quadra poliesportiva) ──────────────────
        # Condicionado ao produto para não sobrescrever com '—' em propostas multi-produto
        if group.productId == 'quadra_poliesportiva':
            _sports: list[str] = []
            if values.get('possui_basquete_adulto'):
                _est = str(values.get('estrutura_basquete_adulto', ''))
                _lbl = _ESTRUTURA_BASQUETE_LABELS.get(_est, _est)
                _sports.append(f'Basquete Adulto ({_lbl})' if _lbl else 'Basquete Adulto')
            if values.get('possui_basquete_juvenil'):
                _sports.append('Basquete Juvenil')
            if values.get('possui_volei'):
                _sports.append('Vôlei')
            if values.get('possui_futebol_futsal'):
                _sports.append('Futebol/Futsal')
            ctx['acessorios_esportivos_descricao'] = f'Acessórios – {", ".join(_sports)}' if _sports else '—'
            ctx['qtde_acessorios_esportivos'] = _fmt_numero(len(_sports)) if _sports else '—'
        ctx['qtde_tela_superior'] = '1,00' if values.get('possui_tela_superior') else '—'
        if not ctx.get('cor_tela_superior'):
            ctx['cor_tela_superior'] = '—'

        # ── EVA: quantidade em metros lineares (perímetro da quadra) ──────
        if values.get('possui_eva'):
            try:
                ctx['qtde_eva'] = _fmt_numero(
                    2 * (float(values.get('largura', 0)) + float(values.get('comprimento', 0)))
                )
            except (TypeError, ValueError):
                ctx['qtde_eva'] = '1,00'
        else:
            ctx['qtde_eva'] = '—'
    return ctx


def _replace_placeholders(slide, context: dict) -> None:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                _replace_in_paragraph(paragraph, context)
        if shape.shape_type == 19:  # TABLE
            for row in shape.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        _replace_in_paragraph(paragraph, context)


def _replace_in_paragraph(paragraph, context: dict) -> None:
    if not paragraph.runs:
        return
    full_text = "".join(run.text for run in paragraph.runs)
    if "{{" not in full_text:
        return
    replaced = full_text
    matched_key = None
    for key, value in context.items():
        new = replaced.replace(f"{{{{ {key} }}}}", str(value) if value is not None else "")
        new = new.replace(f"{{{{{key}}}}}", str(value) if value is not None else "")
        if new != replaced:
            matched_key = key
            replaced = new
    if replaced == full_text:
        return
    paragraph.runs[0].text = replaced
    for run in paragraph.runs[1:]:
        run.text = ""
    if matched_key in _RESET_COLOR_KEYS:
        paragraph.runs[0].font.color.rgb = _PRETO


def _add_placeholder_slide(target: Presentation, label: str) -> None:
    blank = target.slide_layouts[6]
    slide = target.slides.add_slide(blank)
    tx_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1.5))
    tf = tx_box.text_frame
    tf.text = f"[Template pendente: {label}]"
    run = tf.paragraphs[0].runs[0]
    run.font.size = Pt(24)
    run.font.color.rgb = RGBColor(0xCC, 0x00, 0x00)
