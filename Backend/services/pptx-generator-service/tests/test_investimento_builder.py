"""Testes de integração do builder de Investimento — gera PPTX real e inspeciona."""
from io import BytesIO
from types import SimpleNamespace

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from slide_merger import build_presentation


GLOBAL_VALUES = {
    "nome_razao_social": "Cliente Teste",
    "nome_contato": "Fulano",
    "endereco_cliente": "Rua A, 1",
    "local_obra": "SP",
    "telefone": "(11) 9999-0000",
    "email": "f@e.com",
    "numero_proposta": "P-2026-100",
    "data_solicitacao": "2026-05-01",
    "data_envio": "2026-05-20",
}

QUADRA_POLI_ASSOALHO_VALUES = {
    "area_total": 420.0,
    "largura": 18.0,
    "comprimento": 30.0,
    "possui_volei": True,
    "possui_futebol_futsal": True,
    "possui_basquete_adulto": True,
    "estrutura_basquete_adulto": "metalica",
    "possui_basquete_juvenil": False,
    "possui_alambrado": True,
    "possui_iluminacao": True,
    "possui_tela_superior": False,
    "possui_tela_sombreamento": False,
}


def _make_group(product_id, variant_id, values, sumario="x"):
    return SimpleNamespace(
        productId=product_id,
        quantity=1,
        variantId=variant_id,
        values=values,
        sumarioText=sumario,
        investimentoRows=[],
    )


def _make_req(slides, groups):
    return SimpleNamespace(
        slides=slides,
        slideIds=None,
        globalValues=GLOBAL_VALUES,
        productGroups=groups,
    )


def _make_slide(slide_id, template_file, dynamic=None, group_index=None):
    return SimpleNamespace(
        slideId=slide_id,
        templateFile=template_file,
        dynamic=dynamic,
        groupIndex=group_index,
    )


def _open(pptx_bytes: bytes) -> Presentation:
    return Presentation(BytesIO(pptx_bytes))


def _find_table(slide):
    for shape in slide.shapes:
        if shape.shape_type == 19:  # TABLE
            return shape.table
    return None


def _all_text(slide) -> str:
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            parts.append(shape.text_frame.text)
        if shape.shape_type == 19:
            for row in shape.table.rows:
                for cell in row.cells:
                    parts.append(cell.text_frame.text)
    return "\n".join(parts)


def _build_quadra_poli_pptx(variant_id="assoalho", values=None):
    values = values if values is not None else QUADRA_POLI_ASSOALHO_VALUES
    group = _make_group("quadra_poliesportiva", variant_id, values, "Quadra Poliesportiva")
    slide = _make_slide(
        f"investimento_{variant_id}_quadra_poliesportiva",
        "slides/global/investimento_base.pptx",
        dynamic="investimento",
        group_index=0,
    )
    req = _make_req([slide], [group])
    return build_presentation(req)


class TestComposeInvestimento:
    def test_gera_slide_com_tabela(self):
        prs = _open(_build_quadra_poli_pptx())
        assert len(prs.slides) == 1
        table = _find_table(prs.slides[0])
        assert table is not None

    def test_numero_de_linhas_bate_com_catalogo(self):
        # values: piso + volei + futsal + basquete_adulto + alambrado + iluminacao = 6 items
        # → 6 + 2 (header + total) = 8 linhas
        prs = _open(_build_quadra_poli_pptx())
        table = _find_table(prs.slides[0])
        assert len(table.rows) == 8

    def test_header_tem_4_colunas_corretas(self):
        prs = _open(_build_quadra_poli_pptx())
        table = _find_table(prs.slides[0])
        header = [table.cell(0, c).text_frame.text for c in range(4)]
        assert header == ["Descrição", "Qtde.", "Unid.", "Valor"]

    def test_ultima_linha_eh_valor_total(self):
        prs = _open(_build_quadra_poli_pptx())
        table = _find_table(prs.slides[0])
        last = len(table.rows) - 1
        assert "VALOR TOTAL" in table.cell(last, 0).text_frame.text
        assert "R$" in table.cell(last, 3).text_frame.text

    def test_ancora_removida(self):
        prs = _open(_build_quadra_poli_pptx())
        full = _all_text(prs.slides[0])
        assert "{{tabela_investimento}}" not in full

    def test_titulo_presente(self):
        prs = _open(_build_quadra_poli_pptx())
        full = _all_text(prs.slides[0])
        assert "2. Investimento" in full

    def test_qtde_piso_formata_com_virgula(self):
        prs = _open(_build_quadra_poli_pptx())
        table = _find_table(prs.slides[0])
        # piso é a primeira linha de item (row 1)
        assert table.cell(1, 1).text_frame.text == "420,00"
        assert table.cell(1, 2).text_frame.text == "m²"

    def test_acessorio_volei_tem_bold(self):
        prs = _open(_build_quadra_poli_pptx())
        table = _find_table(prs.slides[0])
        # row 2 = vôlei (depois do piso)
        cell = table.cell(2, 0)
        runs = cell.text_frame.paragraphs[0].runs
        bold_texts = [r.text for r in runs if r.font.bold]
        assert any("Vôlei" in t for t in bold_texts)

    def test_piso_assoalho_tem_subtexto_demarcacoes_em_vermelho(self):
        prs = _open(_build_quadra_poli_pptx())
        table = _find_table(prs.slides[0])
        # piso ocupa 2 parágrafos (descrição + "Demarcações inclusas")
        cell = table.cell(1, 0)
        paragraphs = cell.text_frame.paragraphs
        assert len(paragraphs) >= 2
        assert "Demarcações inclusas" in paragraphs[1].text

    def test_placeholders_template_substituidos(self):
        # Apenas o âncora foi removido; placeholders do título/footer não existem aqui.
        # Garante que nenhum placeholder remanescente aparece.
        prs = _open(_build_quadra_poli_pptx())
        full = _all_text(prs.slides[0])
        assert "{{" not in full

    def test_sem_acessorios_so_renderiza_piso_e_total(self):
        values = {
            "area_total": 100.0,
            "possui_volei": False,
            "possui_futebol_futsal": False,
            "possui_basquete_adulto": False,
            "possui_basquete_juvenil": False,
            "possui_alambrado": False,
            "possui_iluminacao": False,
            "possui_tela_superior": False,
            "possui_tela_sombreamento": False,
        }
        prs = _open(_build_quadra_poli_pptx(values=values))
        table = _find_table(prs.slides[0])
        assert len(table.rows) == 3  # header + piso + total

    def test_variant_poliuretano_descricao_com_espessura(self):
        values = {**QUADRA_POLI_ASSOALHO_VALUES, "tipo_poliuretano": "b9"}
        prs = _open(_build_quadra_poli_pptx(variant_id="poliuretano", values=values))
        table = _find_table(prs.slides[0])
        # row 1 é o piso poliuretano
        cell_text = table.cell(1, 0).text_frame.text
        assert "Poliuretano 9mm" in cell_text


class TestEstiloComercial:
    """Validações de tipografia e cores definidas pela equipe comercial."""

    def test_header_usa_poppins_bold_14pt_branco(self):
        prs = _open(_build_quadra_poli_pptx())
        table = _find_table(prs.slides[0])
        run = table.cell(0, 0).text_frame.paragraphs[0].runs[0]
        assert run.font.name == "Poppins"
        assert run.font.bold is True
        assert run.font.size == Pt(14)
        assert run.font.color.rgb == RGBColor(0xFF, 0xFF, 0xFF)

    def test_header_fundo_808080(self):
        prs = _open(_build_quadra_poli_pptx())
        table = _find_table(prs.slides[0])
        fill = table.cell(0, 0).fill
        assert fill.fore_color.rgb == RGBColor(0x80, 0x80, 0x80)

    def test_linhas_de_item_usam_poppins_12pt(self):
        prs = _open(_build_quadra_poli_pptx())
        table = _find_table(prs.slides[0])
        # row 2 = vôlei (qualquer linha de item serve)
        run = table.cell(2, 0).text_frame.paragraphs[0].runs[0]
        assert run.font.name == "Poppins"
        assert run.font.size == Pt(12)
        # texto comum (não destaque) → preto
        run_qtde = table.cell(2, 1).text_frame.paragraphs[0].runs[0]
        assert run_qtde.font.color.rgb == RGBColor(0x00, 0x00, 0x00)

    def test_demarcacoes_inclusas_em_vermelho_ff0000(self):
        prs = _open(_build_quadra_poli_pptx())
        table = _find_table(prs.slides[0])
        # piso = row 1; sub-linha "Demarcações inclusas" está no parágrafo [1]
        paragraphs = table.cell(1, 0).text_frame.paragraphs
        assert "Demarcações inclusas" in paragraphs[1].text
        run = paragraphs[1].runs[0]
        assert run.font.color.rgb == RGBColor(0xFF, 0x00, 0x00)

    def test_valor_total_em_vermelho_ff0000(self):
        prs = _open(_build_quadra_poli_pptx())
        table = _find_table(prs.slides[0])
        last = len(table.rows) - 1
        label_run = table.cell(last, 0).text_frame.paragraphs[0].runs[0]
        money_run = table.cell(last, 3).text_frame.paragraphs[0].runs[0]
        assert label_run.font.color.rgb == RGBColor(0xFF, 0x00, 0x00)
        assert label_run.font.bold is True
        assert money_run.font.color.rgb == RGBColor(0xFF, 0x00, 0x00)


class TestAlturaDoHeader:
    """Header não pode crescer proporcionalmente à quantidade de linhas."""

    _HEADER_HEIGHT_LIMITE = Inches(0.5)

    def _values_com_n_itens(self, n_acessorios: int):
        """Liga `n_acessorios` esportes (até 4) para variar o número de linhas."""
        toggles = [
            "possui_volei",
            "possui_futebol_futsal",
            "possui_basquete_adulto",
            "possui_basquete_juvenil",
        ]
        values = dict(QUADRA_POLI_ASSOALHO_VALUES)
        for i, k in enumerate(toggles):
            values[k] = i < n_acessorios
        values["possui_alambrado"] = False
        values["possui_iluminacao"] = False
        return values

    def test_header_com_poucas_linhas_respeita_limite(self):
        # 3 linhas: header + piso + total
        values = self._values_com_n_itens(0)
        prs = _open(_build_quadra_poli_pptx(values=values))
        table = _find_table(prs.slides[0])
        assert len(table.rows) == 3
        assert table.rows[0].height == self._HEADER_HEIGHT_LIMITE

    def test_header_com_muitas_linhas_mantem_mesma_altura(self):
        values = self._values_com_n_itens(4)
        prs = _open(_build_quadra_poli_pptx(values=values))
        table = _find_table(prs.slides[0])
        assert table.rows[0].height == self._HEADER_HEIGHT_LIMITE

    def test_altura_do_header_eh_constante_independente_da_qtde_de_itens(self):
        alturas = []
        for n in (0, 2, 4):
            values = self._values_com_n_itens(n)
            prs = _open(_build_quadra_poli_pptx(values=values))
            table = _find_table(prs.slides[0])
            alturas.append(table.rows[0].height)
        # Todas as alturas iguais
        assert alturas[0] == alturas[1] == alturas[2]


class TestPlaceholdersDescricao:
    """Resolução de {{ key }} nas descrições do catálogo a partir do ctx."""

    def _values_completo(self):
        return {
            **QUADRA_POLI_ASSOALHO_VALUES,
            "sistema_alambrado": "gaiola",
            "altura_alambrado_fundos": 4.0,
            "altura_alambrado_laterais": 4.0,
            "comprimento_alambrado_laterais": 30.0,
            "comprimento_alambrado_fundos": 18.0,
            "galvanizacao": "fogo",
            "travamento": ["travamento_superior", "travamento_inferior"],
            "quantidade_projetores": 8,
            "potencia_projetores": 200,
            "quantidade_postes_iluminacao": 4,
            "altura_postes_iluminacao": 9,
            "possui_tela_superior": True,
            "cor_tela_superior": "verde",
            "possui_tela_sombreamento": True,
        }

    def _row_text(self, table, row_idx):
        return table.cell(row_idx, 0).text_frame.text

    def test_alambrado_interpola_sistema_galvanizacao_e_travamento(self):
        prs = _open(_build_quadra_poli_pptx(values=self._values_completo()))
        table = _find_table(prs.slides[0])
        # row mapping: piso, volei, futsal, basquete_adulto, alambrado, iluminacao, tela_sup, tela_somb, total
        # alambrado é o 5º item (row 5 = item index 4)
        # Vou achar pelo conteúdo
        textos = [self._row_text(table, i) for i in range(len(table.rows))]
        alambrado = next(t for t in textos if t.startswith("Alambrado"))
        assert "Sistema Gaiola" in alambrado
        assert "tubos galvanizados a fogo" in alambrado
        assert "tela revestida em PVC" in alambrado
        # travamento_descricao deve estar resolvido
        assert "travamento" in alambrado.lower()
        # nenhum placeholder remanescente
        assert "{{" not in alambrado

    def test_iluminacao_interpola_projetores_postes_potencia_altura(self):
        prs = _open(_build_quadra_poli_pptx(values=self._values_completo()))
        table = _find_table(prs.slides[0])
        textos = [self._row_text(table, i) for i in range(len(table.rows))]
        ilum = next(t for t in textos if t.startswith("Iluminação"))
        assert "8 projetores" in ilum
        assert "200W" in ilum
        assert "4 postes" in ilum
        assert "altura de 9m" in ilum
        assert "{{" not in ilum

    def test_tela_superior_interpola_cor(self):
        prs = _open(_build_quadra_poli_pptx(values=self._values_completo()))
        table = _find_table(prs.slides[0])
        textos = [self._row_text(table, i) for i in range(len(table.rows))]
        tela = next(t for t in textos if t.startswith("Tela superior"))
        assert "PVC cor verde" in tela
        assert "{{" not in tela

    def test_tela_sombreamento_sem_placeholders(self):
        prs = _open(_build_quadra_poli_pptx(values=self._values_completo()))
        table = _find_table(prs.slides[0])
        textos = [self._row_text(table, i) for i in range(len(table.rows))]
        tela = next(t for t in textos if t.startswith("Tela de sombreamento"))
        assert "sombreamento 80%" in tela
        assert "{{" not in tela

    def test_chave_ausente_vira_vazio_sem_quebrar(self):
        # Liga iluminação sem informar quantidade_projetores etc.
        values = {
            **QUADRA_POLI_ASSOALHO_VALUES,
            "possui_iluminacao": True,
            "possui_alambrado": False,
            # propositalmente sem quantidade_projetores
        }
        prs = _open(_build_quadra_poli_pptx(values=values))
        table = _find_table(prs.slides[0])
        textos = [table.cell(i, 0).text_frame.text for i in range(len(table.rows))]
        ilum = next(t for t in textos if t.startswith("Iluminação"))
        # context_builder preenche faltantes com '—' (fallback)
        assert "{{" not in ilum


QUADRA_TENIS_VALUES = {
    "largura": 10.97,
    "comprimento": 23.77,
    "area_total": 260.66,
    "possui_playcushion": True,
    "possui_alambrado": True,
    "sistema_alambrado": "gaiola",
    "altura_alambrado_fundos": 4.0,
    "altura_alambrado_laterais": 4.0,
    "comprimento_alambrado_laterais": 23.77,
    "comprimento_alambrado_fundos": 10.97,
    "galvanizacao": "fogo",
    "possui_iluminacao": True,
    "quantidade_projetores": 8,
    "potencia_projetores": 200,
    "quantidade_postes_iluminacao": 4,
    "altura_postes_iluminacao": 9,
    "travamento": ["travamento_superior", "travamento_inferior"],
}


def _build_quadra_tenis_pptx(variant_id, values=None, quantity=1):
    v = values if values is not None else QUADRA_TENIS_VALUES
    group = _make_group("quadra_tenis", variant_id, v, "Quadra de Tênis")
    group.quantity = quantity
    slide = _make_slide(
        f"investimento_{variant_id}_quadra_tenis",
        "slides/global/investimento_base.pptx",
        dynamic="investimento",
        group_index=0,
    )
    req = _make_req([slide], [group])
    return build_presentation(req)


class TestQuadraTenisRenderizacao:
    def test_piso_asfaltico_completo(self):
        prs = _open(_build_quadra_tenis_pptx("piso_asfaltico"))
        table = _find_table(prs.slides[0])
        textos = [table.cell(i, 0).text_frame.text for i in range(len(table.rows))]
        # piso + acessorio + alambrado + iluminacao + playcushion + total + header
        assert len(table.rows) == 7
        assert any("Base Asfáltica" in t for t in textos)
        assert any("Acessório Tênis" in t for t in textos)
        assert any(t.startswith("Alambrado") for t in textos)
        assert any(t.startswith("Iluminação") for t in textos)
        assert any("PlayCushion" in t for t in textos)
        assert any("VALOR TOTAL" in t for t in textos)

    def test_saibro_com_kit(self):
        values = {**QUADRA_TENIS_VALUES, "possui_kit_saibro": True, "possui_playcushion": False}
        prs = _open(_build_quadra_tenis_pptx("saibro", values=values))
        table = _find_table(prs.slides[0])
        textos = [table.cell(i, 0).text_frame.text for i in range(len(table.rows))]
        assert any("Saibro" in t for t in textos)
        assert any("Kit saibro" in t for t in textos)
        assert not any("PlayCushion" in t for t in textos)

    def test_grama_natural_sem_extras(self):
        values = {**QUADRA_TENIS_VALUES, "possui_kit_saibro": True, "possui_playcushion": True}
        prs = _open(_build_quadra_tenis_pptx("grama_natural", values=values))
        table = _find_table(prs.slides[0])
        textos = [table.cell(i, 0).text_frame.text for i in range(len(table.rows))]
        assert any("Grama Natural" in t for t in textos)
        # variant_filter remove kit_saibro e playcushion mesmo com toggles ligados
        assert not any("Kit saibro" in t for t in textos)
        assert not any("PlayCushion" in t for t in textos)

    def test_quantidade_acessorio_tenis_usa_group_quantity(self):
        prs = _open(_build_quadra_tenis_pptx("piso_asfaltico", quantity=2))
        table = _find_table(prs.slides[0])
        # Acessório Tênis é row 2 (após header e piso). Sua qtde deve ser '2,00'.
        for r_idx in range(len(table.rows)):
            desc = table.cell(r_idx, 0).text_frame.text
            if "Acessório Tênis" in desc:
                assert table.cell(r_idx, 1).text_frame.text == "2,00"
                break
        else:
            pytest.fail("Linha 'Acessório Tênis' não encontrada")

    def test_alambrado_interpola_placeholders_em_quadra_tenis(self):
        prs = _open(_build_quadra_tenis_pptx("piso_asfaltico"))
        table = _find_table(prs.slides[0])
        textos = [table.cell(i, 0).text_frame.text for i in range(len(table.rows))]
        alambrado = next(t for t in textos if t.startswith("Alambrado"))
        assert "Sistema Gaiola" in alambrado
        assert "tubos galvanizados a fogo" in alambrado
        assert "{{" not in alambrado

    def test_grama_natural_unidade_M2_maiuscula(self):
        prs = _open(_build_quadra_tenis_pptx("grama_natural"))
        table = _find_table(prs.slides[0])
        # piso é row 1; sua unidade está em col 2 e deve ser "M²" (literal do template antigo)
        assert table.cell(1, 2).text_frame.text == "M²"

    def test_saibro_unidade_acessorio_eh_conj_abreviado(self):
        values = {**QUADRA_TENIS_VALUES, "possui_playcushion": False}
        prs = _open(_build_quadra_tenis_pptx("saibro", values=values))
        table = _find_table(prs.slides[0])
        # Acessório do saibro tem unidade "Conj." (literal)
        for r_idx in range(len(table.rows)):
            desc = table.cell(r_idx, 0).text_frame.text
            if desc.startswith("Acessórios:"):
                assert table.cell(r_idx, 2).text_frame.text == "Conj."
                break
        else:
            pytest.fail("Acessório do saibro não encontrado")


PADEL_VALUES = {
    "area_total": 200.0,
    "largura": 10.0,
    "comprimento": 20.0,
    "possui_alambrado": True,  # padel sempre tem
    "sistema_alambrado": "gaiola",
    "altura_alambrado_fundos": 3.0,
    "altura_alambrado_laterais": 3.0,
    "comprimento_alambrado_laterais": 20.0,
    "comprimento_alambrado_fundos": 10.0,
    "galvanizacao": "fogo",
    "travamento": ["travamento_superior"],
}


def _build_padel_pptx(values=None, quantity=1, possui_acessorio=False):
    v = dict(values if values is not None else PADEL_VALUES)
    v["possui_acessorio_padel"] = possui_acessorio
    group = _make_group("padel", "grama_sintetica", v, "Padel")
    group.quantity = quantity
    slide = _make_slide(
        "investimento_padel",
        "slides/global/investimento_base.pptx",
        dynamic="investimento",
        group_index=0,
    )
    req = _make_req([slide], [group])
    return build_presentation(req)


class TestPadelRenderizacao:
    def test_minimo_2_linhas_de_itens_mais_header_e_total(self):
        prs = _open(_build_padel_pptx())
        table = _find_table(prs.slides[0])
        # piso + estrutura + total + header = 4
        assert len(table.rows) == 4

    def test_com_acessorio_tem_3_linhas_de_itens(self):
        prs = _open(_build_padel_pptx(possui_acessorio=True))
        table = _find_table(prs.slides[0])
        assert len(table.rows) == 5  # header + 3 itens + total

    def test_piso_grama_sintetica_descricao(self):
        prs = _open(_build_padel_pptx())
        table = _find_table(prs.slides[0])
        assert "Grama Sintética Limonta" in table.cell(1, 0).text_frame.text
        assert table.cell(1, 2).text_frame.text == "m²"

    def test_estrutura_vidro_iluminacao(self):
        prs = _open(_build_padel_pptx())
        table = _find_table(prs.slides[0])
        textos = [table.cell(i, 0).text_frame.text for i in range(len(table.rows))]
        estrutura = next(t for t in textos if "Estrutura de Vidro" in t)
        # Texto literal, sem placeholders
        assert "04 postes" in estrutura
        assert "08 projetores led 200W" in estrutura
        assert "{{" not in estrutura

    def test_acessorio_padel_qtde_usa_group_quantity(self):
        prs = _open(_build_padel_pptx(quantity=3, possui_acessorio=True))
        table = _find_table(prs.slides[0])
        for i in range(len(table.rows)):
            desc = table.cell(i, 0).text_frame.text
            if "Suporte da rede" in desc:
                assert table.cell(i, 1).text_frame.text == "3,00"
                assert table.cell(i, 2).text_frame.text == "Unid."
                break
        else:
            pytest.fail("Acessório padel não encontrado")

    def test_estrutura_qtde_usa_area_alambrado(self):
        # alambrado em padel = 20*3 + 10*3 = 90 m²
        prs = _open(_build_padel_pptx())
        table = _find_table(prs.slides[0])
        for i in range(len(table.rows)):
            desc = table.cell(i, 0).text_frame.text
            if "Estrutura de Vidro" in desc:
                assert table.cell(i, 1).text_frame.text == "90,00"
                break
        else:
            pytest.fail("Estrutura não encontrada")


def _build_softplay_pptx(values):
    group = _make_group("softplay", "padrao", values, "Softplay")
    slide = _make_slide(
        "investimento_softplay",
        "slides/global/investimento_base.pptx",
        dynamic="investimento",
        group_index=0,
    )
    req = _make_req([slide], [group])
    return build_presentation(req)


class TestSoftplayRenderizacao:
    def test_softplay_minimalista_tem_3_linhas(self):
        values = {"area_total": 50.0, "espessura_sbr": 3, "espessura_epdm": 1}
        prs = _open(_build_softplay_pptx(values))
        table = _find_table(prs.slides[0])
        assert len(table.rows) == 3  # header + 1 item + total

    def test_softplay_descricao_com_espessura(self):
        values = {"area_total": 50.0, "espessura_sbr": 3, "espessura_epdm": 1}
        prs = _open(_build_softplay_pptx(values))
        table = _find_table(prs.slides[0])
        desc = table.cell(1, 0).text_frame.text
        assert "Playground – Sistema Softplay" in desc
        assert "Moldado 'in loco'" in desc
        assert "Espessura final: 4cm" in desc

    def test_softplay_unidade_m2_e_qtde_area(self):
        values = {"area_total": 75.5, "espessura_sbr": 2, "espessura_epdm": 1}
        prs = _open(_build_softplay_pptx(values))
        table = _find_table(prs.slides[0])
        assert table.cell(1, 1).text_frame.text == "75,50"
        assert table.cell(1, 2).text_frame.text == "m²"


PICKLEBALL_VALUES = {
    "area_total": 80.0,
    "largura": 6.10,
    "comprimento": 13.41,
    "possui_alambrado": True,
    "sistema_alambrado": "gaiola",
    "altura_alambrado_fundos": 3.0,
    "altura_alambrado_laterais": 3.0,
    "comprimento_alambrado_laterais": 13.41,
    "comprimento_alambrado_fundos": 6.10,
    "galvanizacao": "fogo",
    "travamento": ["travamento_superior"],
    "possui_iluminacao": True,
    "quantidade_projetores": 4,
    "potencia_projetores": 150,
    "quantidade_postes_iluminacao": 2,
}


def _build_pickleball_pptx(values=None):
    v = values if values is not None else PICKLEBALL_VALUES
    group = _make_group("pickleball", "padrao", v, "Pickleball")
    slide = _make_slide(
        "investimento_pickleball",
        "slides/global/investimento_base.pptx",
        dynamic="investimento",
        group_index=0,
    )
    req = _make_req([slide], [group])
    return build_presentation(req)


class TestPickleballRenderizacao:
    def test_pickleball_completo_tem_5_linhas(self):
        prs = _open(_build_pickleball_pptx())
        table = _find_table(prs.slides[0])
        assert len(table.rows) == 5  # header + 3 itens + total

    def test_piso_texto_com_acessorio_embutido(self):
        prs = _open(_build_pickleball_pptx())
        table = _find_table(prs.slides[0])
        desc = table.cell(1, 0).text_frame.text
        assert "Quadra de Pickleball" in desc
        assert "Incluso acessório" in desc
        assert table.cell(1, 2).text_frame.text == "M²"

    def test_iluminacao_compacta_sem_altura(self):
        prs = _open(_build_pickleball_pptx())
        table = _find_table(prs.slides[0])
        textos = [table.cell(i, 0).text_frame.text for i in range(len(table.rows))]
        ilum = next(t for t in textos if t.startswith("Iluminação"))
        assert "4 projetores LED 150W" in ilum
        assert "2 postes" in ilum
        # versão compacta: sem "altura de Xm"
        assert "altura" not in ilum

    def test_alambrado_interpola_e_unidade_M2(self):
        prs = _open(_build_pickleball_pptx())
        table = _find_table(prs.slides[0])
        for i in range(len(table.rows)):
            if table.cell(i, 0).text_frame.text.startswith("Alambrado"):
                assert table.cell(i, 2).text_frame.text == "M²"
                assert "Sistema Gaiola" in table.cell(i, 0).text_frame.text
                break
        else:
            pytest.fail("Linha alambrado não encontrada")


BEACH_TENIS_VALUES = {
    "area_total": 78.0,
    "largura": 6.0,
    "comprimento": 13.0,
    "possui_alambrado": True,
    "sistema_alambrado": "gaiola",
    "altura_alambrado_fundos": 3.0,
    "altura_alambrado_laterais": 3.0,
    "comprimento_alambrado_laterais": 13.0,
    "comprimento_alambrado_fundos": 6.0,
    "galvanizacao": "fogo",
    "travamento": ["travamento_superior"],
    "possui_iluminacao": True,
    "quantidade_projetores": 6,
    "potencia_projetores": 200,
    "quantidade_postes_iluminacao": 2,
    "altura_postes_iluminacao": 7,
    "possui_eva": True,
}


def _build_beach_tenis_pptx(values=None, quantity=1):
    v = values if values is not None else BEACH_TENIS_VALUES
    group = _make_group("beach_tenis", "padrao", v, "Beach Tennis")
    group.quantity = quantity
    slide = _make_slide(
        "investimento_beach_tenis",
        "slides/global/investimento_base.pptx",
        dynamic="investimento",
        group_index=0,
    )
    req = _make_req([slide], [group])
    return build_presentation(req)


class TestBeachTenisRenderizacao:
    def test_completo_tem_7_linhas(self):
        prs = _open(_build_beach_tenis_pptx())
        table = _find_table(prs.slides[0])
        # header + 5 itens (piso + alambrado + ilum + acessorio + eva) + total
        assert len(table.rows) == 7

    def test_piso_textos_cbt(self):
        prs = _open(_build_beach_tenis_pptx())
        table = _find_table(prs.slides[0])
        desc = table.cell(1, 0).text_frame.text
        assert "Beach Tennis" in desc
        assert "drenagem" in desc
        assert "areia de quartzo" in desc

    def test_iluminacao_sem_prefixo_iluminacao(self):
        # No template antigo o item de iluminação NÃO começa com "Iluminação – "
        prs = _open(_build_beach_tenis_pptx())
        table = _find_table(prs.slides[0])
        textos = [table.cell(i, 0).text_frame.text for i in range(len(table.rows))]
        # Linha de iluminação começa com "6 projetores..."
        ilum = next(t for t in textos if t.startswith("6 projetores"))
        assert "200W" in ilum
        assert "altura de 7m" in ilum

    def test_eva_qtde_eh_perimetro(self):
        # Perímetro = 2 * (6 + 13) = 38
        prs = _open(_build_beach_tenis_pptx())
        table = _find_table(prs.slides[0])
        for i in range(len(table.rows)):
            if "Proteção EVA" in table.cell(i, 0).text_frame.text:
                assert table.cell(i, 1).text_frame.text == "38,00"
                assert table.cell(i, 2).text_frame.text == "ML"
                break
        else:
            pytest.fail("Proteção EVA não encontrada")

    def test_acessorio_qtde_usa_group_quantity(self):
        prs = _open(_build_beach_tenis_pptx(quantity=3))
        table = _find_table(prs.slides[0])
        for i in range(len(table.rows)):
            desc = table.cell(i, 0).text_frame.text
            if desc.startswith("Acessórios com/sem"):
                assert table.cell(i, 1).text_frame.text == "3,00"
                break
        else:
            pytest.fail("Acessório beach tênis não encontrado")

    def test_sem_eva_quando_desligado(self):
        values = {**BEACH_TENIS_VALUES, "possui_eva": False}
        prs = _open(_build_beach_tenis_pptx(values=values))
        table = _find_table(prs.slides[0])
        textos = [table.cell(i, 0).text_frame.text for i in range(len(table.rows))]
        assert not any("Proteção EVA" in t for t in textos)
