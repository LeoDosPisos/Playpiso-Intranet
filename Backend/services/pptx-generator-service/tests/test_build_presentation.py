"""Integration tests for build_presentation.

These tests use the real .pptx template files on disk (no mocking).
Run with: pytest tests/ -v
"""
from io import BytesIO
from pathlib import Path
from types import SimpleNamespace

import pytest
from pptx import Presentation

from slide_merger import build_presentation
from context_builder import _build_base_context, _build_group_context
from placeholder_engine import _replace_placeholders

_SLIDES_DIR = Path(__file__).resolve().parent.parent / "slides"


# ── helpers ───────────────────────────────────────────────────────────────────

GLOBAL_VALUES = {
    "nome_razao_social": "Condomínio Teste",
    "nome_contato": "João Silva",
    "endereco_cliente": "Rua Teste, 123",
    "local_obra": "São Paulo, SP",
    "telefone": "(11) 99999-0000",
    "email": "joao@teste.com",
    "numero_proposta": "P-2026-001",
    "data_solicitacao": "2026-05-01",
    "data_envio": "2026-05-20",
}

QUADRA_TENIS_VALUES = {
    "largura": 10.97,
    "comprimento": 23.77,
    "area_total": 260.66,
    "possui_playcushion": True,
    "possui_alambrado": True,
    "sistema_alambrado": "gaiola",
    "altura_alambrado_fundos": 4.0,
    "galvanizacao": "fogo",
    "possui_iluminacao": False,
    "travamento": ["travamento_superior", "travamento_inferior"],
    "comprimento_alambrado_laterais": 23.77,
    "altura_alambrado_laterais": 4.0,
    "comprimento_alambrado_fundos": 10.97,
}

BEACH_TENIS_VALUES = {
    "largura": 6.0,
    "comprimento": 13.0,
    "area_total": 78.0,
    "possui_alambrado": False,
    "possui_iluminacao": False,
    "possui_eva": True,
    "possui_tela_superior": False,
}


def _make_group(product_id, values, sumario_text="1 produto", quantity=1, variant_id="piso_asfaltico"):
    return SimpleNamespace(
        productId=product_id,
        quantity=quantity,
        variantId=variant_id,
        values=values,
        sumarioText=sumario_text,
        investimentoRows=[],
    )


def _make_slide(slide_id, template_file, dynamic=None, group_index=None):
    return SimpleNamespace(
        slideId=slide_id,
        templateFile=template_file,
        dynamic=dynamic,
        groupIndex=group_index,
    )


def _make_req(slides, global_values=None, product_groups=None):
    return SimpleNamespace(
        slides=slides,
        slideIds=None,
        globalValues=global_values or GLOBAL_VALUES,
        productGroups=product_groups or [],
    )


def _open_pptx(pptx_bytes: bytes) -> Presentation:
    return Presentation(BytesIO(pptx_bytes))


def _unsubstituted_placeholders(prs: Presentation) -> list[str]:
    """Return list of text snippets that still contain {{ }} (recurse em GROUP shapes)."""
    found = []

    def _scan(shapes):
        for shape in shapes:
            if shape.shape_type == 6:  # GROUP
                _scan(shape.shapes)
                continue
            if shape.has_text_frame:
                txt = shape.text_frame.text
                if "{{" in txt:
                    found.append(txt.strip()[:80])
            if shape.shape_type == 19:
                for row in shape.table.rows:
                    for cell in row.cells:
                        txt = cell.text_frame.text
                        if "{{" in txt:
                            found.append(txt.strip()[:80])

    for slide in prs.slides:
        _scan(slide.shapes)
    return found


# ── context unit tests ────────────────────────────────────────────────────────

class TestBuildBaseContext:
    def test_global_fields_present(self):
        group = _make_group("quadra_tenis", QUADRA_TENIS_VALUES, "Quadra de Tênis")
        ctx = _build_base_context(GLOBAL_VALUES, [group])
        assert ctx["nome_razao_social"] == "Condomínio Teste"
        assert ctx["data_solicitacao"] == "01/05/2026"
        assert ctx["data_envio"] == "20/05/2026"

    def test_short_aliases_present(self):
        group = _make_group("quadra_tenis", QUADRA_TENIS_VALUES, "Quadra de Tênis")
        ctx = _build_base_context(GLOBAL_VALUES, [group])
        assert ctx["np"] == "P-2026-001"
        assert ctx["ds"] == "01/05/2026"
        assert ctx["de"] == "20/05/2026"

    def test_sumario_single_product(self):
        group = _make_group("quadra_tenis", QUADRA_TENIS_VALUES, "Quadra de Tênis")
        ctx = _build_base_context(GLOBAL_VALUES, [group])
        assert ctx["sumario"] == "Quadra de Tênis"

    def test_sumario_multi_product(self):
        g1 = _make_group("quadra_tenis", QUADRA_TENIS_VALUES, "Quadra de Tênis")
        g2 = _make_group("beach_tenis", BEACH_TENIS_VALUES, "Beach Tennis")
        ctx = _build_base_context(GLOBAL_VALUES, [g1, g2])
        assert ctx["sumario"] == "1. Quadra de Tênis\n2. Beach Tennis"


class TestBuildGroupContext:
    def test_travamento_descricao(self):
        group = _make_group("quadra_tenis", QUADRA_TENIS_VALUES)
        ctx = _build_group_context(group)
        assert "superior" in ctx["travamento_descricao"]
        assert "inferior" in ctx["travamento_descricao"]

    def test_alambrado_gaiola(self):
        group = _make_group("quadra_tenis", QUADRA_TENIS_VALUES)
        ctx = _build_group_context(group)
        assert (
            "Fundos e laterais com 4,00m de altura de tela em sistema gaiola"
            in ctx["alambrado_descricao"]
        )

    def test_alambrado_trapezio(self):
        values = {
            **QUADRA_TENIS_VALUES,
            "sistema_alambrado": "trapezio",
            "altura_alambrado_fundos": 4.0,
            "altura_alambrado_laterais": 4.0,
        }
        ctx = _build_group_context(_make_group("quadra_tenis", values))
        assert (
            "Fundos com 4,00m de altura de tela e laterais de 4,00m"
            " em sistema trapézio com corrimão de 1,00m de altura"
            in ctx["alambrado_descricao"]
        )

    def test_tela_malha_por_produto(self):
        # Poliesportiva e campo usam 3" fio 12; os demais (e fallback) usam 2" fio 14.
        # (espaço entre "fio" e o número é NBSP,  )
        for pid in ("quadra_poliesportiva", "campo"):
            ctx = _build_group_context(_make_group(pid, QUADRA_TENIS_VALUES))
            assert ctx["tela_malha"] == "3” x 3” fio 12"
        for pid in ("quadra_tenis", "pickleball", "beach_tenis", "padel"):
            ctx = _build_group_context(_make_group(pid, QUADRA_TENIS_VALUES))
            assert ctx["tela_malha"] == "2” x 2” fio 14"

    def test_altura_postes_iluminacao_fmt(self):
        # Altura dos postes (slide secao_iluminacao) formatada "X,XXm" a partir do valor do form.
        ctx = _build_group_context(_make_group("quadra_tenis", {"altura_postes_iluminacao": 8}))
        assert ctx["altura_postes_iluminacao_fmt"] == "8,00m"
        ctx = _build_group_context(_make_group("campo", {"altura_postes_iluminacao": 6}))
        assert ctx["altura_postes_iluminacao_fmt"] == "6,00m"
        ctx = _build_group_context(_make_group("quadra_tenis", {"altura_postes_iluminacao": 7.5}))
        assert ctx["altura_postes_iluminacao_fmt"] == "7,50m"   # vírgula, 2 casas
        ctx = _build_group_context(_make_group("quadra_tenis", {}))
        assert ctx["altura_postes_iluminacao_fmt"] == "—"        # ausente

    def test_placeholders_dentro_de_group_shape(self):
        # O slide padel/secao_iluminacao.pptx tem o texto dentro de um GROUP shape;
        # o motor deve recursar e substituir os placeholders aninhados.
        ctx = _build_group_context(_make_group(
            "padel",
            {"quantidade_postes_iluminacao": 4, "quantidade_projetores": 16, "potencia_projetores": 400},
            variant_id="grama_sintetica",
        ))
        prs = Presentation(str(_SLIDES_DIR / "padel" / "secao_iluminacao.pptx"))
        for slide in prs.slides:
            _replace_placeholders(slide, ctx)
        assert _unsubstituted_placeholders(prs) == []

        def _group_text(shapes):
            for shp in shapes:
                if shp.shape_type == 6:
                    yield from _group_text(shp.shapes)
                elif shp.has_text_frame:
                    yield shp.text_frame.text

        joined = "\n".join(t for slide in prs.slides for t in _group_text(slide.shapes))
        assert "4 postes de iluminação semicirculares" in joined

    def test_alambrado_cores_formatadas(self):
        values = {
            **QUADRA_TENIS_VALUES,
            "cor_tubo_alambrado": "branco",
            "cor_tela_malha_alambrado": "preto",
            "cor_tela_alambrado": "verde",
        }
        ctx = _build_group_context(_make_group("quadra_tenis", values))
        assert ctx["cor_tubo"] == "branca"        # "pintados na cor branca" (fem.)
        assert ctx["cor_tela_malha"] == "preto"   # "PVC preto" (masc.)
        assert "com tela na cor verde" in ctx["alambrado_descricao"]

    def test_alambrado_cores_fallback(self):
        ctx = _build_group_context(_make_group("quadra_tenis", QUADRA_TENIS_VALUES))
        assert ctx["cor_tubo"] == "—"
        assert ctx["cor_tela_malha"] == "—"
        assert "com tela na cor" not in ctx["alambrado_descricao"]

    def test_galvanizacao_mapped(self):
        group = _make_group("quadra_tenis", QUADRA_TENIS_VALUES)
        ctx = _build_group_context(group)
        assert ctx["galvanizacao"] == "a fogo"

    def test_area_alambrado_computed(self):
        group = _make_group("quadra_tenis", QUADRA_TENIS_VALUES)
        ctx = _build_group_context(group)
        # 23.77*4.0 + 10.97*4.0 = 95.08 + 43.88 = 138.96
        assert ctx["area_alambrado"] != "—"

    def test_no_alambrado_returns_dash(self):
        values = {**QUADRA_TENIS_VALUES, "possui_alambrado": False}
        group = _make_group("quadra_tenis", values)
        ctx = _build_group_context(group)
        assert ctx["area_alambrado"] == "—"

    def test_alambrado_especial(self):
        values = {
            **QUADRA_TENIS_VALUES,
            "sistema_alambrado": "especial",
            "comprimento_alambrado_lateral_1": 23.77,
            "altura_alambrado_lateral_1": 4.0,
            "espacamento_postes_tubos_lateral_1": 2.5,
            "comprimento_alambrado_lateral_2": 23.77,
            "altura_alambrado_lateral_2": 3.5,
            "espacamento_postes_tubos_lateral_2": 2.5,
            "comprimento_alambrado_fundo_1": 10.97,
            "altura_alambrado_fundo_1": 4.0,
            "espacamento_postes_tubos_fundo_1": 2.0,
            "comprimento_alambrado_fundo_2": 10.97,
            "altura_alambrado_fundo_2": 3.0,
            "espacamento_postes_tubos_fundo_2": 2.0,
        }
        group = _make_group("quadra_tenis", values)
        ctx = _build_group_context(group)

        assert ctx["sistema_alambrado"] == "Especial"
        descricao = ctx["alambrado_descricao"].lower()
        assert "especial" in descricao
        assert "lateral 1" in descricao
        assert "fundo 2" in descricao
        # área = 23.77*4 + 23.77*3.5 + 10.97*4 + 10.97*3 = 95.08+83.20+43.88+32.91 = 255.07
        assert ctx["area_alambrado"] != "—"

    def test_eva_qtde_computed(self):
        group = _make_group("beach_tenis", BEACH_TENIS_VALUES)
        ctx = _build_group_context(group)
        # perímetro = 2*(6+13) = 38
        assert ctx["qtde_eva"] != "—"

    def test_espessura_total_softplay(self):
        group = _make_group("softplay", {"espessura_sbr": 3, "espessura_epdm": 1})
        ctx = _build_group_context(group)
        assert ctx["espessura_total"] == "4"

    def test_espessura_total_softplay_decimal(self):
        group = _make_group("softplay", {"espessura_sbr": 3, "espessura_epdm": 1.5})
        ctx = _build_group_context(group)
        assert ctx["espessura_total"] == "4,5"

    def test_contexts_are_isolated(self):
        """Each call to _build_group_context returns independent data."""
        g1 = _make_group("quadra_tenis", QUADRA_TENIS_VALUES)
        g2 = _make_group("beach_tenis", {**BEACH_TENIS_VALUES, "possui_alambrado": True,
                                          "sistema_alambrado": "trapezio",
                                          "altura_alambrado_fundos": 3.0,
                                          "altura_alambrado_laterais": 3.0})
        ctx1 = _build_group_context(g1)
        ctx2 = _build_group_context(g2)

        assert ctx1["sistema_alambrado"] != ctx2["sistema_alambrado"]
        assert ctx1["alambrado_descricao"] != ctx2["alambrado_descricao"]


# ── integration tests (require real .pptx files) ─────────────────────────────

class TestBuildPresentation:
    def test_single_slide_capa(self):
        group = _make_group("quadra_tenis", QUADRA_TENIS_VALUES, "Quadra de Tênis")
        slides = [_make_slide("capa", "slides/global/capa.pptx")]
        req = _make_req(slides, product_groups=[group])
        pptx_bytes = build_presentation(req)
        prs = _open_pptx(pptx_bytes)
        assert len(prs.slides) == 1

    def test_dados_cliente_placeholders_substituted(self):
        group = _make_group("quadra_tenis", QUADRA_TENIS_VALUES, "Quadra de Tênis")
        slides = [_make_slide("dados_cliente", "slides/global/dados_cliente.pptx")]
        req = _make_req(slides, product_groups=[group])
        pptx_bytes = build_presentation(req)
        prs = _open_pptx(pptx_bytes)
        remaining = _unsubstituted_placeholders(prs)
        assert remaining == [], f"Placeholders não substituídos: {remaining}"

    def test_no_slides_returns_empty_pptx(self):
        group = _make_group("quadra_tenis", QUADRA_TENIS_VALUES)
        req = _make_req([], product_groups=[group])
        pptx_bytes = build_presentation(req)
        prs = _open_pptx(pptx_bytes)
        assert len(prs.slides) == 0

    def test_investimento_quadra_tenis_placeholders(self):
        group = _make_group("quadra_tenis", QUADRA_TENIS_VALUES, "Quadra de Tênis")
        slides = [
            _make_slide(
                "investimento_piso_asfaltico_quadra_tenis",
                "slides/quadra_tenis/investimento_piso_asfaltico.pptx",
                group_index=0,
            )
        ]
        req = _make_req(slides, product_groups=[group])
        pptx_bytes = build_presentation(req)
        prs = _open_pptx(pptx_bytes)
        assert len(prs.slides) >= 1
        remaining = _unsubstituted_placeholders(prs)
        assert remaining == [], f"Placeholders não substituídos: {remaining}"

    def test_fechamentos_quadra_tenis_created(self):
        group = _make_group("quadra_tenis", QUADRA_TENIS_VALUES, "Quadra de Tênis")
        slides = [
            _make_slide(
                "fechamentos_quadra_tenis",
                "slides/_comum/fechamentos_base.pptx",
                dynamic="fechamentos",
                group_index=0,
            )
        ]
        req = _make_req(slides, product_groups=[group])
        pptx_bytes = build_presentation(req)
        prs = _open_pptx(pptx_bytes)
        # possui_alambrado=True, possui_iluminacao=False → 1 seção ativa → 1 slide
        assert len(prs.slides) == 1

    def test_multi_produto_contextos_isolados(self):
        """Com 2 grupos, cada slide deve usar o contexto do seu groupIndex."""
        g1 = _make_group(
            "quadra_tenis",
            {**QUADRA_TENIS_VALUES, "possui_alambrado": True, "sistema_alambrado": "gaiola"},
            "Quadra de Tênis",
        )
        g2 = _make_group(
            "beach_tenis",
            {**BEACH_TENIS_VALUES, "possui_alambrado": False},
            "Beach Tennis",
        )
        slides = [
            _make_slide("capa", "slides/global/capa.pptx"),
            _make_slide("dados_cliente", "slides/global/dados_cliente.pptx"),
            _make_slide(
                "investimento_piso_asfaltico_quadra_tenis",
                "slides/quadra_tenis/investimento_piso_asfaltico.pptx",
                group_index=0,
            ),
            _make_slide(
                "investimento_beach_tenis",
                "slides/beach_tenis/investimento.pptx",
                group_index=1,
            ),
        ]
        req = _make_req(slides, product_groups=[g1, g2])
        pptx_bytes = build_presentation(req)
        prs = _open_pptx(pptx_bytes)
        assert len(prs.slides) == 4
        remaining = _unsubstituted_placeholders(prs)
        assert remaining == [], f"Placeholders não substituídos: {remaining}"

    def test_img_counter_does_not_bleed_between_requests(self):
        """Duas chamadas independentes não devem acumular o contador de imagens."""
        group = _make_group("quadra_tenis", QUADRA_TENIS_VALUES, "Quadra de Tênis")
        slides = [_make_slide("capa", "slides/global/capa.pptx")]
        req = _make_req(slides, product_groups=[group])

        bytes1 = build_presentation(req)
        bytes2 = build_presentation(req)
        # Ambas devem gerar arquivos válidos com o mesmo número de slides
        assert len(_open_pptx(bytes1).slides) == len(_open_pptx(bytes2).slides)


# ── Bloco opcional de "projetos realizados" por produto ──────────────────────
class TestComposeProjetos:
    """compose_projetos: lê slides/<product_subdir>/projetos/*.pptx; pasta ausente/vazia = no-op."""

    @staticmethod
    def _make_minimal_pptx(path):
        p = Presentation()
        p.slides.add_slide(p.slide_layouts[6])  # layout em branco
        p.save(str(path))

    @staticmethod
    def _fresh_merged():
        from slide_registry import SLIDE_H, SLIDE_W
        merged = Presentation()
        merged.slide_width = SLIDE_W
        merged.slide_height = SLIDE_H
        return merged

    def test_sem_nada_no_op(self, monkeypatch, tmp_path):
        from dynamic_composer import compose_projetos
        monkeypatch.setattr("slide_registry.SLIDES_DIR", str(tmp_path))
        merged = self._fresh_merged()
        compose_projetos(merged, "quadra_tenis", "piso_asfaltico", {}, [0])
        assert len(merged.slides) == 0

    def test_variante_especifica_tem_precedencia(self, monkeypatch, tmp_path):
        """Variante com .pptx é usada; nível produto NÃO entra (sem concatenação)."""
        from dynamic_composer import compose_projetos
        monkeypatch.setattr("slide_registry.SLIDES_DIR", str(tmp_path))
        base = tmp_path / "quadra_tenis" / "projetos"
        base.mkdir(parents=True)
        self._make_minimal_pptx(base / "comum.pptx")               # nível produto
        (base / "piso_asfaltico").mkdir()
        self._make_minimal_pptx(base / "piso_asfaltico" / "v1.pptx")
        self._make_minimal_pptx(base / "piso_asfaltico" / "v2.pptx")
        merged = self._fresh_merged()
        compose_projetos(merged, "quadra_tenis", "piso_asfaltico", {}, [0])
        assert len(merged.slides) == 2  # só a variante; nível produto NÃO concatena

    def test_fallback_para_nivel_produto_quando_variante_vazia(self, monkeypatch, tmp_path):
        from dynamic_composer import compose_projetos
        monkeypatch.setattr("slide_registry.SLIDES_DIR", str(tmp_path))
        base = tmp_path / "quadra_tenis" / "projetos"
        base.mkdir(parents=True)
        (base / "saibro").mkdir()  # variante existe mas sem .pptx
        self._make_minimal_pptx(base / "padrao.pptx")
        merged = self._fresh_merged()
        compose_projetos(merged, "quadra_tenis", "saibro", {}, [0])
        assert len(merged.slides) == 1  # cai p/ nível produto

    def test_fallback_para_nivel_produto_quando_variante_ausente(self, monkeypatch, tmp_path):
        from dynamic_composer import compose_projetos
        monkeypatch.setattr("slide_registry.SLIDES_DIR", str(tmp_path))
        base = tmp_path / "padel" / "projetos"
        base.mkdir(parents=True)
        self._make_minimal_pptx(base / "comum.pptx")
        merged = self._fresh_merged()
        # Sem subpasta 'grama_sintetica/' → cai p/ slides/padel/projetos/.
        compose_projetos(merged, "padel", "grama_sintetica", {}, [0])
        assert len(merged.slides) == 1

    def test_arquivos_adicionados_em_ordem_alfabetica(self, monkeypatch, tmp_path):
        from dynamic_composer import compose_projetos
        monkeypatch.setattr("slide_registry.SLIDES_DIR", str(tmp_path))
        variant_dir = tmp_path / "quadra_tenis" / "projetos" / "piso_asfaltico"
        variant_dir.mkdir(parents=True)
        self._make_minimal_pptx(variant_dir / "02_projeto.pptx")
        self._make_minimal_pptx(variant_dir / "01_projeto.pptx")
        merged = self._fresh_merged()
        compose_projetos(merged, "quadra_tenis", "piso_asfaltico", {}, [0])
        assert len(merged.slides) == 2  # 2 arquivos, 1 slide cada

    def test_usa_subdir_mapeado_quadra_poli(self, monkeypatch, tmp_path):
        """quadra_poliesportiva → 'quadra_poli' via _PRODUCT_SLIDES_DIR."""
        from dynamic_composer import compose_projetos
        monkeypatch.setattr("slide_registry.SLIDES_DIR", str(tmp_path))
        variant_dir = tmp_path / "quadra_poli" / "projetos" / "assoalho"
        variant_dir.mkdir(parents=True)
        self._make_minimal_pptx(variant_dir / "01.pptx")
        merged = self._fresh_merged()
        compose_projetos(merged, "quadra_poliesportiva", "assoalho", {}, [0])
        assert len(merged.slides) == 1

    def test_variant_id_vazio_pula_para_nivel_produto(self, monkeypatch, tmp_path):
        """Sem variantId definido, vai direto ao nível produto."""
        from dynamic_composer import compose_projetos
        monkeypatch.setattr("slide_registry.SLIDES_DIR", str(tmp_path))
        base = tmp_path / "softplay" / "projetos"
        base.mkdir(parents=True)
        self._make_minimal_pptx(base / "01.pptx")
        merged = self._fresh_merged()
        compose_projetos(merged, "softplay", None, {}, [0])
        assert len(merged.slides) == 1

    def test_build_presentation_aceita_dynamic_projetos(self):
        """Smoke: build_presentation aceita slide com dynamic='projetos' sem crashar.

        Robusto: a real `slides/quadra_tenis/projetos/<variant>/` pode existir e
        adicionar N slides; o que importa é que o branch não quebra e o deck
        contém pelo menos capa + encerramento.
        """
        group = _make_group("quadra_tenis", QUADRA_TENIS_VALUES, "Quadra de Tênis")
        slides = [
            _make_slide("capa", "slides/global/capa.pptx"),
            _make_slide(
                "projetos_quadra_tenis__piso_asfaltico",
                "", dynamic="projetos", group_index=0,
            ),
            _make_slide("encerramento", "slides/global/encerramento.pptx"),
        ]
        req = _make_req(slides, product_groups=[group])
        prs = _open_pptx(build_presentation(req))
        assert len(prs.slides) >= 2
