import logging
import os
import time
from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from context_builder import _build_base_context, _build_group_context, _build_context, _is_truthy
from dynamic_composer import compose_fechamentos, compose_acessorios
from investimento import compose_investimento
from placeholder_engine import _replace_placeholders
from slide_copier import _copy_slide
from slide_registry import SLIDES_DIR, SLIDE_W, SLIDE_H, _slide_template_path, _resolve_slide_path

_INVESTIMENTO_BASE_REL = "global/investimento_base.pptx"

logger = logging.getLogger("pptx_generator.merger")

# Mapeia slideId → [(título_da_seção, chave_no_contexto), ...]
# Seções cujo valor de contexto seja falsy são removidas antes da substituição.
CONDITIONAL_SECTIONS: dict[str, list[tuple[str, str]]] = {}


def _remove_section(slide, section_title: str) -> bool:
    """Remove título + conteúdo de uma seção identificada pelo texto exato do título."""
    spTree = slide.shapes._spTree

    title_el = None
    title_top = None
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip() == section_title:
            title_el = shape.element
            title_top = shape.top
            break

    if title_el is None:
        logger.warning("section_not_found", extra={"section": section_title})
        return False

    content_el = None
    min_dist = float('inf')
    for shape in slide.shapes:
        if shape.element is title_el:
            continue
        if shape.has_text_frame and shape.top > title_top:
            dist = shape.top - title_top
            if dist < min_dist:
                min_dist = dist
                content_el = shape.element

    spTree.remove(title_el)
    if content_el is not None:
        spTree.remove(content_el)

    logger.info("section_removed", extra={"section": section_title})
    return True


def _add_from_file_with_replacement(
    merged: Presentation, file_path: str, label: str, context: dict, img_counter: list[int]
) -> None:
    src = Presentation(file_path)
    for src_slide in src.slides:
        new_slide = _copy_slide(merged, src_slide, img_counter)

        if label in CONDITIONAL_SECTIONS:
            for section_title, context_key in CONDITIONAL_SECTIONS[label]:
                if not _is_truthy(context.get(context_key)):
                    _remove_section(new_slide, section_title)

        _replace_placeholders(new_slide, context)

    logger.info("slide_loaded", extra={"label": label, "slide_count": len(src.slides)})


def _add_placeholder_slide(target: Presentation, label: str) -> None:
    blank = target.slide_layouts[6]
    slide = target.slides.add_slide(blank)
    tx_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1.5))
    tf = tx_box.text_frame
    tf.text = f"[Template pendente: {label}]"
    run = tf.paragraphs[0].runs[0]
    run.font.size = Pt(24)
    run.font.color.rgb = RGBColor(0xCC, 0x00, 0x00)


def build_presentation(req) -> bytes:
    t0 = time.perf_counter()

    use_rich_slides = hasattr(req, 'slides') and req.slides is not None
    base_ctx = _build_base_context(req.globalValues, req.productGroups)
    group_ctxs = [_build_group_context(g) for g in req.productGroups]
    single_group = len(req.productGroups) == 1
    groups_by_index = {i: g for i, g in enumerate(req.productGroups)}
    img_counter = [0]

    def _ctx_for(group_idx):
        if group_idx is not None and group_idx < len(group_ctxs):
            return {**base_ctx, **group_ctxs[group_idx]}
        if single_group:
            return {**base_ctx, **group_ctxs[0]}
        return base_ctx

    merged = Presentation()
    merged.slide_width = SLIDE_W
    merged.slide_height = SLIDE_H

    if use_rich_slides:
        for slide_entry in req.slides:
            slide_id = slide_entry.slideId
            dynamic   = slide_entry.dynamic
            group_idx = slide_entry.groupIndex
            ctx = _ctx_for(group_idx)

            if dynamic == "fechamentos":
                product_id = slide_id.removeprefix("fechamentos_")
                group  = groups_by_index.get(group_idx or 0)
                values = dict(group.values) if group else {}
                base_rel  = slide_entry.templateFile.removeprefix("slides/")
                base_path = _resolve_slide_path(base_rel)
                compose_fechamentos(merged, base_path, product_id, values, ctx, img_counter)
            elif dynamic == "acessorios":
                product_id = slide_id.removeprefix("acessorios_")
                group  = groups_by_index.get(group_idx or 0)
                values = dict(group.values) if group else {}
                base_rel  = slide_entry.templateFile.removeprefix("slides/")
                base_path = _resolve_slide_path(base_rel)
                compose_acessorios(merged, base_path, product_id, values, ctx, img_counter)
            elif dynamic == "investimento":
                group  = groups_by_index.get(group_idx or 0)
                if group is None:
                    _add_placeholder_slide(merged, slide_id)
                    continue
                values = dict(group.values)
                # group.quantity é metadado do grupo, não do form values — injetar
                # com prefixo _ para itens do catálogo que usam quantity como qtde
                # (ex: Acessório Tênis da quadra_tenis).
                values["_quantity"] = group.quantity
                base_path = os.path.join(SLIDES_DIR, _INVESTIMENTO_BASE_REL)
                try:
                    compose_investimento(
                        merged, base_path, group.productId, group.variantId,
                        values, ctx, img_counter,
                    )
                except NotImplementedError:
                    # Produto ainda não migrado para o catálogo dinâmico → cai no template legado.
                    base_rel  = slide_entry.templateFile.removeprefix("slides/")
                    file_path = os.path.join(SLIDES_DIR, base_rel)
                    if os.path.exists(file_path):
                        _add_from_file_with_replacement(merged, file_path, slide_id, ctx, img_counter)
                    else:
                        _add_placeholder_slide(merged, slide_id)
            else:
                base_rel  = slide_entry.templateFile.removeprefix("slides/")
                file_path = os.path.join(SLIDES_DIR, base_rel)
                if os.path.exists(file_path):
                    _add_from_file_with_replacement(merged, file_path, slide_id, ctx, img_counter)
                else:
                    _add_placeholder_slide(merged, slide_id)
    else:
        legacy_ctx = _build_context(req.globalValues, req.productGroups)
        for slide_id in (req.slideIds or []):
            file_path = _slide_template_path(slide_id)
            if file_path:
                _add_from_file_with_replacement(merged, file_path, slide_id, legacy_ctx, img_counter)
            else:
                _add_placeholder_slide(merged, slide_id)

    output = BytesIO()
    merged.save(output)
    result = output.getvalue()

    logger.info(
        "presentation_built",
        extra={
            "slide_count": len(merged.slides),
            "size_bytes": len(result),
            "duration_ms": round((time.perf_counter() - t0) * 1000, 1),
        },
    )
    return result
