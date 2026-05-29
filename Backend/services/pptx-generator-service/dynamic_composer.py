import copy
import logging
import os

from pptx import Presentation

from context_builder import _is_truthy
from placeholder_engine import _replace_placeholders
from slide_copier import _copy_slide, _shift_shape_top
from slide_registry import _PRODUCT_SLIDES_DIR, _resolve_slide_path

logger = logging.getLogger("pptx_generator.merger")

_CM_TO_EMU                  = 914400 / 2.54
_CONTENT_TOP_EMU            = round(0.5 * _CM_TO_EMU)
_ACESSORIOS_CONTENT_TOP_EMU = round(2.25 * _CM_TO_EMU)
_SECTION_GAP_EMU            = round(0 * _CM_TO_EMU)
_ACESSORIOS_SECTION_GAP_EMU = round(-0.1 * _CM_TO_EMU)

_FECHAMENTOS_SECTIONS: list[tuple[str, str]] = [
    ("possui_alambrado",         "alambrado"),
    ("possui_iluminacao",        "iluminacao"),
    ("possui_tela_superior",     "tela_superior"),
    ("possui_tela_sombreamento", "tela_sombreamento"),
]


def _get_active_acessorios_sections(product_id: str, values: dict) -> list[str]:
    if product_id == "padel":
        return ["padel"] if _is_truthy(values.get("possui_acessorio_padel")) else []

    if product_id == "pickleball":
        return ["pickleball"] if _is_truthy(values.get("possui_rede_pickleball")) else []

    if product_id == "quadra_tenis":
        sections: list[str] = []
        if _is_truthy(values.get("incluir_rede_tenis")):
            sections.append("rede_tenis")
        if _is_truthy(values.get("possui_kit_saibro")):
            sections.append("kit_saibro")
        return sections

    sections: list[str] = []

    if _is_truthy(values.get("possui_basquete_adulto")):
        estrutura = str(values.get("estrutura_basquete_adulto") or "comum")
        if estrutura in ("metalica", "hidraulica", "comum"):
            sections.append(f"basquete_adulto_{estrutura}")
        else:
            sections.append("basquete_adulto_comum")

    if _is_truthy(values.get("possui_basquete_juvenil")):
        sections.append("basquete_juvenil")

    if _is_truthy(values.get("possui_volei")):
        sections.append("volei")

    if _is_truthy(values.get("possui_tenis")):
        sections.append("tenis")

    if _is_truthy(values.get("possui_futebol_futsal")):
        tipo = str(values.get("tipo_futsal") or "padrao")
        sections.append(f"futsal_{tipo}" if tipo in ("padrao", "mini_trave") else "futsal_padrao")

    return sections


def compose_fechamentos(
    merged: Presentation,
    base_path: str,
    product_id: str,
    values: dict,
    ctx: dict,
    img_counter: list[int],
) -> None:
    active = [
        name for field, name in _FECHAMENTOS_SECTIONS if _is_truthy(values.get(field))
    ]

    # Padel: alambrado é obrigatório, renderiza o slide base mesmo sem nenhuma
    # das seções dinâmicas ativas (as informações de alambrado já estão no base).
    if not active:
        if product_id == "padel":
            if not os.path.exists(base_path):
                logger.warning("fechamentos_base não encontrado: %s", base_path)
                return
            src_base = Presentation(base_path)
            new_slide = _copy_slide(merged, src_base.slides[0], img_counter)
            _replace_placeholders(new_slide, ctx)
        return

    pages = [active] if len(active) <= 3 else [active[:2], active[2:]]
    product_subdir = _PRODUCT_SLIDES_DIR.get(product_id, product_id)

    for page_sections in pages:
        if not os.path.exists(base_path):
            logger.warning("fechamentos_base não encontrado: %s", base_path)
            return

        src_base = Presentation(base_path)
        new_slide = _copy_slide(merged, src_base.slides[0], img_counter)

        sp_tree = new_slide.shapes._spTree
        current_top = _CONTENT_TOP_EMU
        gap = 0 if len(page_sections) == 2 else _SECTION_GAP_EMU

        for section_name in page_sections:
            comp_path = _resolve_slide_path(os.path.join(product_subdir, f"secao_{section_name}.pptx"))
            if not os.path.exists(comp_path):
                logger.warning("componente não encontrado: %s", comp_path)
                continue
            src_comp = Presentation(comp_path)
            comp_slide = src_comp.slides[0]

            comp_height = max(
                (shape.top + shape.height for shape in comp_slide.shapes),
                default=0,
            )

            for shape in comp_slide.shapes:
                el = copy.deepcopy(shape._element)
                _shift_shape_top(el, current_top)
                sp_tree.append(el)

            current_top += comp_height + gap

        _replace_placeholders(new_slide, ctx)


def _list_projetos_pptx(directory: str) -> list[str]:
    """Lista .pptx ordenados alfabeticamente em ``directory``; [] se ausente/vazio.

    Ignora arquivos ocultos (``.``) e subpastas — só arquivos .pptx no nível do diretório.
    """
    if not os.path.isdir(directory):
        return []
    return sorted(
        f for f in os.listdir(directory)
        if f.lower().endswith(".pptx")
        and not f.startswith(".")
        and os.path.isfile(os.path.join(directory, f))
    )


def compose_projetos(
    merged: Presentation,
    product_id: str,
    variant_id: str | None,
    ctx: dict,
    img_counter: list[int],
) -> None:
    """Anexa os slides estáticos de 'projetos realizados' por (produto, variante).

    Convenção com fallback:
      1) Tenta ``slides/<product_subdir>/projetos/<variant_id>/*.pptx``.
      2) Se vazio/ausente, cai para ``slides/<product_subdir>/projetos/*.pptx``
         (nível produto, comum a todas as variantes).
      3) Sem nenhum dos dois → bloco silenciosamente omitido (log info).

    Arquivos ordenados alfabeticamente; cada um pode conter 1+ slides; todos
    são copiados como slides independentes (diferente de compose_fechamentos/
    acessorios, que empilham seções numa única base).
    """
    product_subdir = _PRODUCT_SLIDES_DIR.get(product_id, product_id)

    # 1) Variante específica.
    if variant_id:
        variant_dir = _resolve_slide_path(os.path.join(product_subdir, "projetos", str(variant_id)))
        files = _list_projetos_pptx(variant_dir)
        if files:
            for fn in files:
                src = Presentation(os.path.join(variant_dir, fn))
                for slide in src.slides:
                    new_slide = _copy_slide(merged, slide, img_counter)
                    _replace_placeholders(new_slide, ctx)
            return

    # 2) Fallback nível produto.
    product_dir = _resolve_slide_path(os.path.join(product_subdir, "projetos"))
    files = _list_projetos_pptx(product_dir)
    if not files:
        logger.info(
            "projetos: nada para %s/%s (variante=%s, produto=%s)",
            product_id, variant_id or "—", variant_dir if variant_id else "—", product_dir,
        )
        return

    for fn in files:
        src = Presentation(os.path.join(product_dir, fn))
        for slide in src.slides:
            new_slide = _copy_slide(merged, slide, img_counter)
            _replace_placeholders(new_slide, ctx)


def compose_acessorios(
    merged: Presentation,
    base_path: str,
    product_id: str,
    values: dict,
    ctx: dict,
    img_counter: list[int],
) -> None:
    active = _get_active_acessorios_sections(product_id, values)
    if not active:
        return

    product_subdir = _PRODUCT_SLIDES_DIR.get(product_id, product_id)
    pages = [active[i:i + 4] for i in range(0, len(active), 4)]

    for page_sections in pages:
        if not os.path.exists(base_path):
            logger.warning("acessorios_base não encontrado: %s", base_path)
            return

        src_base = Presentation(base_path)
        new_slide = _copy_slide(merged, src_base.slides[0], img_counter)
        sp_tree = new_slide.shapes._spTree
        current_top = _ACESSORIOS_CONTENT_TOP_EMU

        for section_name in page_sections:
            comp_path = _resolve_slide_path(
                os.path.join(product_subdir, f"secao_acessorio_{section_name}.pptx")
            )
            if not os.path.exists(comp_path):
                logger.warning("componente acessório não encontrado: %s", comp_path)
                continue

            src_comp = Presentation(comp_path)
            comp_slide = src_comp.slides[0]
            comp_height = max(
                (shape.top + shape.height for shape in comp_slide.shapes),
                default=0,
            )
            for shape in comp_slide.shapes:
                el = copy.deepcopy(shape._element)
                _shift_shape_top(el, current_top)
                sp_tree.append(el)
            current_top += comp_height + _ACESSORIOS_SECTION_GAP_EMU

        _replace_placeholders(new_slide, ctx)
