import copy
import os

from pptx import Presentation
from pptx.parts.image import ImagePart
from pptx.opc.packuri import PackURI

_R_EMBED = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'
_NS_P    = 'http://schemas.openxmlformats.org/presentationml/2006/main'


def _shift_shape_top(el, delta_emu: int) -> None:
    from pptx.oxml.ns import qn
    for tag in (qn("p:spPr"), qn("p:grpSpPr")):
        spPr = el.find(tag)
        if spPr is None:
            continue
        xfrm = spPr.find(qn("a:xfrm"))
        if xfrm is None:
            continue
        off = xfrm.find(qn("a:off"))
        if off is not None:
            off.set("y", str(int(off.get("y", "0")) + delta_emu))
        return


def _copy_background(src_slide, dest_slide) -> None:
    src_cSld = src_slide._element.find(f'{{{_NS_P}}}cSld')
    src_bg = src_cSld.find(f'{{{_NS_P}}}bg') if src_cSld is not None else None

    if src_bg is None:
        return

    dest_cSld = dest_slide._element.find(f'{{{_NS_P}}}cSld')
    existing_bg = dest_cSld.find(f'{{{_NS_P}}}bg')
    if existing_bg is not None:
        dest_cSld.remove(existing_bg)

    dest_cSld.insert(0, copy.deepcopy(src_bg))


def _copy_slide(merged: Presentation, src_slide, img_counter: list[int]):
    blank_layout = merged.slide_layouts[6]
    new_slide = merged.slides.add_slide(blank_layout)

    rId_map = {}
    for rId, rel in src_slide.part.rels.items():
        if '/image' in rel.reltype:
            img_counter[0] += 1
            src_img = rel.target_part
            ext = os.path.splitext(str(src_img.partname))[1]
            new_partname = PackURI(f'/ppt/media/m{img_counter[0]}{ext}')
            new_img = ImagePart(new_partname, src_img.content_type, merged.part.package, src_img.blob)
            new_rId = new_slide.part.relate_to(new_img, rel.reltype)
            rId_map[rId] = new_rId

    spTree_copy = copy.deepcopy(src_slide.shapes._spTree)
    for el in spTree_copy.iter():
        old = el.get(_R_EMBED)
        if old and old in rId_map:
            el.set(_R_EMBED, rId_map[old])

    sp = new_slide.shapes._spTree
    sp.clear()
    for el in spTree_copy:
        sp.append(el)

    _copy_background(src_slide, new_slide)

    return new_slide
