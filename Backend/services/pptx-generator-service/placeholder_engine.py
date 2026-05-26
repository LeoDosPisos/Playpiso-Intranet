import re

from pptx.dml.color import RGBColor

_PRETO = RGBColor(0x1C, 0x1C, 0x1C)
_RESET_COLOR_KEYS = frozenset({"sumario"})

# \s matches NBSP (\xa0), tab, regular space etc., tolerating whitespace artifacts
# that PowerPoint sometimes injects when editing placeholder text.
_PLACEHOLDER_RE = re.compile(r"\{\{\s*([A-Za-z_][A-Za-z0-9_]*)\s*\}\}")


def _replace_placeholders(slide, context: dict) -> None:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                _replace_in_paragraph(paragraph, context)
        if shape.shape_type == 19:  # TABLE
            for row in shape.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        _replace_in_paragraph(paragraph, context)


def _replace_in_paragraph(paragraph, context: dict) -> None:
    if not paragraph.runs:
        return
    full_text = "".join(run.text for run in paragraph.runs)
    if "{{" not in full_text:
        return

    matched_key: str | None = None

    def _sub(match: re.Match) -> str:
        nonlocal matched_key
        key = match.group(1)
        if key in context:
            matched_key = key
            value = context[key]
            return str(value) if value is not None else ""
        return match.group(0)

    replaced = _PLACEHOLDER_RE.sub(_sub, full_text)
    if replaced == full_text:
        return
    paragraph.runs[0].text = replaced
    for run in paragraph.runs[1:]:
        run.text = ""
    if matched_key in _RESET_COLOR_KEYS:
        paragraph.runs[0].font.color.rgb = _PRETO
