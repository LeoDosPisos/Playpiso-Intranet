"""Insere o token {{ altura_postes_iluminacao_fmt }} no slide secao_iluminacao.pptx.

Operação:
 - Localiza o parágrafo dos POSTES (único que contém "altura livre"; o parágrafo das
   cruzetas também referencia {{ quantidade_postes_iluminacao }}, então não serve como
   discriminador).
 - Dentro desse parágrafo, identifica o run que contém o substring
   "altura livre de 8,00m" e troca **só o texto desse run** por
   "altura livre de {{ altura_postes_iluminacao_fmt }}". Nenhum outro run é tocado, então
   toda a formatação (fonte, tamanho, cor, idioma, dirty flag, etc.) é preservada.
 - A frase das cruzetas (com {{ cor_cruzetas }}) fica intocada.

Idempotente: se o template já contiver "{{ altura_postes_iluminacao_fmt }}", o script é
no-op e devolve com sucesso.

Espelha o padrão de scripts/edit_secao_iluminacao_pptx.py (que insere {{ cor_cruzetas }}).

Uso: python scripts/edit_secao_iluminacao_altura_pptx.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation

REPO_ROOT = Path(__file__).resolve().parents[1]
PPTX_PATH = (
    REPO_ROOT
    / "Backend/services/pptx-generator-service/slides/_comum/secao_iluminacao.pptx"
)

# Discriminador do parágrafo dos POSTES (único que menciona altura livre).
POSTES_DISCRIMINATOR = "altura livre"
# Marcador da frase das CRUZETAS, usado só para checagem pós-save (não deve mudar).
CRUZETAS_TOKEN = "{{ quantidade_cruzetas }}"

NEW_TOKEN = "{{ altura_postes_iluminacao_fmt }}"
SEARCH = "altura livre de 8,00m"
REPLACEMENT = f"altura livre de {NEW_TOKEN}"


def _para_text(paragraph) -> str:
    return "".join(run.text for run in paragraph.runs)


def _iter_text_paragraphs(prs):
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                yield paragraph


def main() -> int:
    if not PPTX_PATH.exists():
        raise SystemExit(f"Template não encontrado: {PPTX_PATH}")

    prs = Presentation(str(PPTX_PATH))

    postes_paragraphs = [
        p for p in _iter_text_paragraphs(prs) if POSTES_DISCRIMINATOR in _para_text(p)
    ]
    if not postes_paragraphs:
        raise SystemExit(
            f"Não encontrei parágrafo com '{POSTES_DISCRIMINATOR}'. "
            "O template foi alterado e o script precisa ser revisto."
        )
    if len(postes_paragraphs) > 1:
        raise SystemExit(
            f"Encontrei {len(postes_paragraphs)} parágrafos com '{POSTES_DISCRIMINATOR}'; "
            "esperava exatamente 1. Aborta para evitar edição ambígua."
        )

    target = postes_paragraphs[0]
    full_text = _para_text(target)

    if NEW_TOKEN in full_text:
        print(f"[no-op] '{NEW_TOKEN}' já presente no parágrafo dos postes — nada a fazer.")
        return 0

    if SEARCH not in full_text:
        raise SystemExit(
            f"O parágrafo dos postes não contém '{SEARCH}'. Texto atual:\n  {full_text!r}"
        )

    # Estratégia preservadora de formatação: localizar o run cujo .text contém
    # 'altura livre de 8,00m' inteiro e fazer replace apenas nesse run.
    target_runs = [r for r in target.runs if SEARCH in r.text]
    if len(target_runs) != 1:
        raise SystemExit(
            f"'{SEARCH}' não está inteiro em um único run do parágrafo dos postes "
            f"(encontrei {len(target_runs)} runs candidatos). "
            "Provavelmente o template foi reformatado; revise o script antes de prosseguir."
        )

    run = target_runs[0]
    original_run_text = run.text
    run.text = original_run_text.replace(SEARCH, REPLACEMENT, 1)

    prs.save(str(PPTX_PATH))
    print(f"[ok] Atualizei {PPTX_PATH.relative_to(REPO_ROOT)}")
    print(f"     Run antes:  {original_run_text!r}")
    print(f"     Run depois: {run.text!r}")

    # Verificações pós-save
    prs2 = Presentation(str(PPTX_PATH))
    paragraphs2 = list(_iter_text_paragraphs(prs2))
    full_doc = "\n".join(_para_text(p) for p in paragraphs2)

    occurrences = full_doc.count(NEW_TOKEN)
    if occurrences != 1:
        raise SystemExit(
            f"Verificação falhou: '{NEW_TOKEN}' aparece {occurrences}× no slide; esperava 1."
        )

    cruzetas_paragraphs = [p for p in paragraphs2 if CRUZETAS_TOKEN in _para_text(p)]
    if not cruzetas_paragraphs:
        raise SystemExit(
            "Verificação falhou: parágrafo das cruzetas não foi encontrado após o save."
        )

    postes_after = [p for p in paragraphs2 if POSTES_DISCRIMINATOR in _para_text(p)][0]
    print(f"[ok] Verificações pós-save: '{NEW_TOKEN}' presente 1×; parágrafo das cruzetas intacto.")
    print(f"     Postes: {len(postes_after.runs)} run(s) preservado(s).")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
