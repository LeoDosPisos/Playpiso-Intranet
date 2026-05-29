"""Divide o slide acessorio.pptx (Beach Tennis) em dois templates condicionais.

Origem (atualmente referenciado pelo slideRegistry):
    Backend/services/pptx-generator-service/slides/beach_tenis/acessorio.pptx

Esse arquivo contém DOIS blocos no mesmo slide:
 - "Sem  regulagem" (shape no canto superior do conteúdo)
 - "Com  regulagem" (shape no canto inferior)

Saída:
 - acessorio_sem_regulagem.pptx: cópia do original sem o bloco "Com regulagem".
 - acessorio_com_regulagem.pptx: cópia sem o bloco "Sem regulagem", com o
   bloco "Com regulagem" REPOSICIONADO para o top/left original do
   bloco "Sem regulagem" (evita um buraco grande no topo do slide).

Por fim, **apaga** o acessorio.pptx original.

Idempotente: se os dois arquivos novos já existem com o conteúdo certo e o
original não existe, o script termina como no-op.

Uso: python scripts/split_acessorio_beach_tenis_pptx.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation

REPO_ROOT = Path(__file__).resolve().parents[1]
SLIDES_DIR = REPO_ROOT / "Backend/services/pptx-generator-service/slides/beach_tenis"
ORIGINAL = SLIDES_DIR / "acessorio.pptx"
OUT_SEM = SLIDES_DIR / "acessorio_sem_regulagem.pptx"
OUT_COM = SLIDES_DIR / "acessorio_com_regulagem.pptx"

# Texto de âncora — note os DOIS espaços, exatamente como o template gravou.
SEM_ANCHOR = "Sem  regulagem"
COM_ANCHOR = "Com  regulagem"


def _shape_text(shape) -> str:
    if not shape.has_text_frame:
        return ""
    return "\n".join("".join(r.text for r in p.runs) for p in shape.text_frame.paragraphs)


def _find_shape_by_anchor(slide, anchor: str):
    """Acha o único shape cujo texto contém o anchor (ou erra)."""
    matches = [sh for sh in slide.shapes if anchor in _shape_text(sh)]
    if len(matches) != 1:
        raise SystemExit(
            f"Esperava 1 shape contendo {anchor!r}; achei {len(matches)}. "
            "O template pode ter sido alterado; revise o script."
        )
    return matches[0]


def _remove_shape(shape) -> None:
    el = shape._element
    el.getparent().remove(el)


def _build_sem_regulagem() -> None:
    prs = Presentation(str(ORIGINAL))
    slide = prs.slides[0]
    com_shape = _find_shape_by_anchor(slide, COM_ANCHOR)
    _remove_shape(com_shape)
    # Sanidade: confere que ainda tem o bloco SEM e não tem mais o COM.
    texts = [_shape_text(s) for s in slide.shapes]
    if not any(SEM_ANCHOR in t for t in texts):
        raise SystemExit("acessorio_sem_regulagem.pptx: bloco 'Sem regulagem' não encontrado pós-remoção.")
    if any(COM_ANCHOR in t for t in texts):
        raise SystemExit("acessorio_sem_regulagem.pptx: bloco 'Com regulagem' ainda presente pós-remoção.")
    prs.save(str(OUT_SEM))


def _build_com_regulagem() -> None:
    prs = Presentation(str(ORIGINAL))
    slide = prs.slides[0]
    sem_shape = _find_shape_by_anchor(slide, SEM_ANCHOR)
    com_shape = _find_shape_by_anchor(slide, COM_ANCHOR)

    # Guarda a posição original do bloco "Sem regulagem" para reaproveitar.
    sem_top, sem_left = sem_shape.top, sem_shape.left

    _remove_shape(sem_shape)

    # Move o bloco "Com regulagem" para o topo (onde estava o "Sem regulagem").
    com_shape.top = sem_top
    com_shape.left = sem_left

    texts = [_shape_text(s) for s in slide.shapes]
    if not any(COM_ANCHOR in t for t in texts):
        raise SystemExit("acessorio_com_regulagem.pptx: bloco 'Com regulagem' não encontrado pós-remoção.")
    if any(SEM_ANCHOR in t for t in texts):
        raise SystemExit("acessorio_com_regulagem.pptx: bloco 'Sem regulagem' ainda presente pós-remoção.")
    prs.save(str(OUT_COM))


def _output_already_valid(path: Path, must_contain: str, must_not_contain: str) -> bool:
    if not path.exists():
        return False
    try:
        prs = Presentation(str(path))
    except Exception:
        return False
    slide = prs.slides[0]
    texts = [_shape_text(s) for s in slide.shapes]
    has_required = any(must_contain in t for t in texts)
    has_forbidden = any(must_not_contain in t for t in texts)
    return has_required and not has_forbidden


def main() -> int:
    # Idempotência: se os dois saídos já estão válidos e o original já sumiu, no-op.
    if (
        not ORIGINAL.exists()
        and _output_already_valid(OUT_SEM, SEM_ANCHOR, COM_ANCHOR)
        and _output_already_valid(OUT_COM, COM_ANCHOR, SEM_ANCHOR)
    ):
        print("[no-op] Split já aplicado anteriormente; nada a fazer.")
        return 0

    if not ORIGINAL.exists():
        raise SystemExit(
            f"Original não encontrado em {ORIGINAL.relative_to(REPO_ROOT)} e os outputs "
            "esperados estão incompletos. Restaure o arquivo (git checkout) e rode de novo."
        )

    _build_sem_regulagem()
    print(f"[ok] Criado {OUT_SEM.relative_to(REPO_ROOT)}")

    _build_com_regulagem()
    print(f"[ok] Criado {OUT_COM.relative_to(REPO_ROOT)}")

    ORIGINAL.unlink()
    print(f"[ok] Apagado {ORIGINAL.relative_to(REPO_ROOT)}")

    # Verificações finais (releitura)
    if not _output_already_valid(OUT_SEM, SEM_ANCHOR, COM_ANCHOR):
        raise SystemExit("Verificação pós-save FALHOU em acessorio_sem_regulagem.pptx")
    if not _output_already_valid(OUT_COM, COM_ANCHOR, SEM_ANCHOR):
        raise SystemExit("Verificação pós-save FALHOU em acessorio_com_regulagem.pptx")

    print("[ok] Verificações pós-save concluídas com sucesso.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
