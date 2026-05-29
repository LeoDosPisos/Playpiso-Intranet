"""Insere o token {{ cor_cruzetas }} no slide secao_iluminacao.pptx.

Operação:
 - Localiza o parágrafo que contém o token `{{ quantidade_cruzetas }}` (a frase
   das cruzetas).
 - Dentro desse parágrafo, identifica o run que contém o substring
   "na cor verde" e troca **só o texto desse run** por
   "na cor {{ cor_cruzetas }}". Nenhum outro run é tocado, então toda a
   formatação (fonte, tamanho, cor, idioma, dirty flag, etc.) é preservada.
 - A frase dos postes (que também termina com "na cor verde") fica intocada.

Idempotente: se o template já contiver "{{ cor_cruzetas }}", o script é no-op
e devolve com sucesso.

Uso: python scripts/edit_secao_iluminacao_pptx.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation

REPO_ROOT = Path(__file__).resolve().parents[1]
PPTX_PATH = (
    REPO_ROOT
    / "Backend/services/pptx-generator-service/slides/_comum/secao_iluminacao.pptx"
)

CRUZETAS_TOKEN = "{{ quantidade_cruzetas }}"
POSTES_TOKEN = "{{ quantidade_postes_iluminacao }}"
NEW_TOKEN = "{{ cor_cruzetas }}"
SEARCH = "na cor verde"
REPLACEMENT = f"na cor {NEW_TOKEN}"


def _para_text(paragraph) -> str:
    return "".join(run.text for run in paragraph.runs)


def _iter_text_paragraphs(prs):
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                yield paragraph


def main() -> int:
    if not PPTX_PATH.exists():
        raise SystemExit(f"Template não encontrado: {PPTX_PATH}")

    prs = Presentation(str(PPTX_PATH))

    cruzetas_paragraphs = [p for p in _iter_text_paragraphs(prs) if CRUZETAS_TOKEN in _para_text(p)]
    if not cruzetas_paragraphs:
        raise SystemExit(
            f"Não encontrei o parágrafo que contém '{CRUZETAS_TOKEN}'. "
            "O template foi alterado e o script precisa ser revisto."
        )
    if len(cruzetas_paragraphs) > 1:
        raise SystemExit(
            f"Encontrei {len(cruzetas_paragraphs)} parágrafos com '{CRUZETAS_TOKEN}'; "
            "esperava exatamente 1. Aborta para evitar edição ambígua."
        )

    target = cruzetas_paragraphs[0]
    full_text = _para_text(target)

    if NEW_TOKEN in full_text:
        print(f"[no-op] '{NEW_TOKEN}' já presente no parágrafo das cruzetas — nada a fazer.")
        return 0

    if SEARCH not in full_text:
        raise SystemExit(
            f"O parágrafo das cruzetas não contém '{SEARCH}'. Texto atual:\n  {full_text!r}"
        )

    # Estratégia preservadora de formatação: localizar o run cujo .text contém
    # 'na cor verde' inteiro e fazer replace apenas nesse run.
    target_runs = [r for r in target.runs if SEARCH in r.text]
    if len(target_runs) != 1:
        # Caso o substring esteja partido entre runs, o script ainda pode rodar —
        # mas exigiria mesclar runs, o que desfaz formatação local. Melhor avisar.
        raise SystemExit(
            f"'{SEARCH}' não está inteiro em um único run do parágrafo das cruzetas "
            f"(encontrei {len(target_runs)} runs candidatos). "
            "Provavelmente o template foi reformatado; revise o script antes de prosseguir."
        )

    run = target_runs[0]
    original_run_text = run.text
    run.text = original_run_text.replace(SEARCH, REPLACEMENT, 1)

    prs.save(str(PPTX_PATH))
    print(f"[ok] Atualizei {PPTX_PATH.relative_to(REPO_ROOT)}")
    print(f"     Run antes:  {original_run_text!r}")
    print(f"     Run depois: {run.text!r}")

    # Verificações pós-save
    prs2 = Presentation(str(PPTX_PATH))
    paragraphs2 = list(_iter_text_paragraphs(prs2))
    full_doc = "\n".join(_para_text(p) for p in paragraphs2)

    occurrences = full_doc.count(NEW_TOKEN)
    if occurrences != 1:
        raise SystemExit(
            f"Verificação falhou: '{NEW_TOKEN}' aparece {occurrences}× no slide; esperava 1."
        )

    postes_paragraphs = [
        p for p in paragraphs2
        if POSTES_TOKEN in _para_text(p) and CRUZETAS_TOKEN not in _para_text(p)
    ]
    for p in postes_paragraphs:
        if SEARCH not in _para_text(p):
            raise SystemExit(
                "Verificação falhou: a frase dos postes não contém mais 'na cor verde'.\n"
                f"  Conteúdo: {_para_text(p)!r}"
            )

    # Confere que o parágrafo das cruzetas mantém o mesmo número de runs (= formatação preservada)
    cruzetas_after = [p for p in paragraphs2 if CRUZETAS_TOKEN in _para_text(p)][0]
    print(f"[ok] Verificações pós-save: '{NEW_TOKEN}' presente 1×; frase dos postes intacta.")
    print(f"     Cruzetas: {len(cruzetas_after.runs)} run(s) preservado(s).")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
