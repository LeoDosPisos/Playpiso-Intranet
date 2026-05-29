"""Insere o placeholder {{ nome_da_obra }} no slide global/dados_cliente.pptx.

Operação:
 - Localiza, na tabela do slide, a célula cujo `text_frame` contém o placeholder
   `{{ local_obra }}`. É o bloco "Local da obra / {{ local_obra }}" (Row 3, Col 0
   no template atual).
 - Localiza o **parágrafo de valor** dentro dessa célula (o que contém
   `{{ local_obra }}` — *não* o label "Local da obra").
 - Clona o XML desse parágrafo, troca o texto do run para `{{ nome_da_obra }}`
   e insere o clone como **novo primeiro parágrafo** da célula (antes do label
   "Local da obra"). Resultado visual: <nome da obra> acima de "Local da obra".
 - O clone preserva fonte/tamanho/cor/negrito do run original (Poppins Bold ~15pt)
   porque deepcopy do `<a:p>` mantém todas as propriedades.

Idempotente: se o template já contiver `{{ nome_da_obra }}` em qualquer lugar do
slide, o script é no-op e devolve com sucesso.

Espelha o padrão de scripts/edit_secao_iluminacao_pptx.py e
edit_secao_iluminacao_altura_pptx.py.

Uso: python scripts/edit_dados_cliente_nome_da_obra_pptx.py
"""
from __future__ import annotations

import copy
from pathlib import Path

from pptx import Presentation

REPO_ROOT = Path(__file__).resolve().parents[1]
PPTX_PATH = (
    REPO_ROOT
    / "Backend/services/pptx-generator-service/slides/global/dados_cliente.pptx"
)

EXISTING_TOKEN = "{{ local_obra }}"
NEW_TOKEN = "{{ nome_da_obra }}"
EXISTING_LABEL = "Local da obra"


def _para_text(paragraph) -> str:
    return "".join(run.text for run in paragraph.runs)


def _tf_text(text_frame) -> str:
    return "\n".join(_para_text(p) for p in text_frame.paragraphs)


def _iter_table_cells(prs):
    """Itera todas as células de todas as tabelas de todos os slides."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            for row in shape.table.rows:
                for cell in row.cells:
                    yield cell


def _slide_full_text(prs) -> str:
    """Concatena o texto de todos os shapes (incluindo células de tabela)."""
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(_tf_text(shape.text_frame))
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        parts.append(_tf_text(cell.text_frame))
    return "\n".join(parts)


def main() -> int:
    if not PPTX_PATH.exists():
        raise SystemExit(f"Template não encontrado: {PPTX_PATH}")

    prs = Presentation(str(PPTX_PATH))

    # Idempotência: se o token novo já está em qualquer lugar do slide, no-op.
    if NEW_TOKEN in _slide_full_text(prs):
        print(f"[no-op] '{NEW_TOKEN}' já presente em {PPTX_PATH.relative_to(REPO_ROOT)} — nada a fazer.")
        return 0

    # Localiza a célula com {{ local_obra }} (esperado: exatamente 1).
    target_cells = [c for c in _iter_table_cells(prs) if EXISTING_TOKEN in _tf_text(c.text_frame)]
    if not target_cells:
        raise SystemExit(
            f"Não encontrei célula com '{EXISTING_TOKEN}'. "
            "O template foi alterado e o script precisa ser revisto."
        )
    if len(target_cells) > 1:
        raise SystemExit(
            f"Encontrei {len(target_cells)} células com '{EXISTING_TOKEN}'; esperava exatamente 1."
        )

    cell = target_cells[0]
    tf = cell.text_frame
    paragraphs = list(tf.paragraphs)

    # Sanity check: a célula deve conter o label "Local da obra".
    if not any(EXISTING_LABEL in _para_text(p) for p in paragraphs):
        raise SystemExit(
            f"A célula com '{EXISTING_TOKEN}' não tem o label '{EXISTING_LABEL}'. "
            "Estrutura inesperada do template — aborta para evitar edição ambígua."
        )

    # Acha o parágrafo de VALOR (o que tem o placeholder, via texto concatenado dos
    # runs — o template típico parte `{{ local_obra }}` em 3 runs `{{ ` / `local_obra` / ` }}`).
    value_paragraphs = [p for p in paragraphs if EXISTING_TOKEN in _para_text(p)]
    if len(value_paragraphs) != 1:
        raise SystemExit(
            f"Esperava 1 parágrafo com '{EXISTING_TOKEN}' na célula, encontrei {len(value_paragraphs)}."
        )
    value_p = value_paragraphs[0]

    # Clona o parágrafo inteiro (preserva paragraph properties + run properties).
    new_p_el = copy.deepcopy(value_p._p)

    # Estratégia: concatena o texto de todos os runs do clone, faz o replace,
    # joga o texto resultante no PRIMEIRO run e remove os demais. Como o placeholder
    # costuma estar partido em 3 runs com formatação IDÊNTICA, mesclá-los no primeiro
    # não perde nada visualmente — e a formatação do primeiro run é justamente a do
    # placeholder (que é o que queremos para o novo token).
    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    r_qname = f"{{{a_ns}}}r"
    t_qname = f"{{{a_ns}}}t"
    r_elements = new_p_el.findall(r_qname)
    if not r_elements:
        raise SystemExit("Falha interna: parágrafo de valor clonado não tem nenhum run.")

    def _r_text(r_el):
        t_el = r_el.find(t_qname)
        return (t_el.text or "") if t_el is not None else ""

    current_text = "".join(_r_text(r) for r in r_elements)
    if EXISTING_TOKEN not in current_text:
        raise SystemExit(
            f"Texto concatenado do parágrafo de valor não contém '{EXISTING_TOKEN}'. "
            f"Texto atual: {current_text!r}"
        )
    new_text = current_text.replace(EXISTING_TOKEN, NEW_TOKEN, 1)

    first_r = r_elements[0]
    first_t = first_r.find(t_qname)
    if first_t is None:
        raise SystemExit("Falha interna: primeiro run do clone não tem elemento <a:t>.")
    original_text = first_t.text
    first_t.text = new_text
    for extra_r in r_elements[1:]:
        new_p_el.remove(extra_r)

    # Insere o novo parágrafo ANTES do primeiro parágrafo existente da célula.
    first_p_el = paragraphs[0]._p
    first_p_el.addprevious(new_p_el)

    prs.save(str(PPTX_PATH))
    print(f"[ok] Atualizei {PPTX_PATH.relative_to(REPO_ROOT)}")
    print(f"     Texto concat. antes: {current_text!r} (primeiro run era {original_text!r}, {len(r_elements)} runs)")
    print(f"     Texto do novo parágrafo: {first_t.text!r} (1 run após merge)")

    # Verificações pós-save
    prs2 = Presentation(str(PPTX_PATH))
    full = _slide_full_text(prs2)

    n_new = full.count(NEW_TOKEN)
    if n_new != 1:
        raise SystemExit(
            f"Verificação falhou: '{NEW_TOKEN}' aparece {n_new}× no slide; esperava 1."
        )
    n_local = full.count(EXISTING_TOKEN)
    if n_local != 1:
        raise SystemExit(
            f"Verificação falhou: '{EXISTING_TOKEN}' aparece {n_local}× no slide; esperava 1."
        )
    if EXISTING_LABEL not in full:
        raise SystemExit(
            f"Verificação falhou: label '{EXISTING_LABEL}' sumiu após o save."
        )

    # Confirma a ORDEM: na célula alvo (já reaberta), o parágrafo do novo token
    # deve vir antes do label "Local da obra".
    target_cells2 = [
        c for c in _iter_table_cells(prs2) if NEW_TOKEN in _tf_text(c.text_frame)
    ]
    if len(target_cells2) != 1:
        raise SystemExit(
            f"Verificação falhou: esperava 1 célula com '{NEW_TOKEN}' após save, achei {len(target_cells2)}."
        )
    cell2 = target_cells2[0]
    paragraphs2 = list(cell2.text_frame.paragraphs)
    idx_new = next(
        (i for i, p in enumerate(paragraphs2) if NEW_TOKEN in _para_text(p)), -1
    )
    idx_label = next(
        (i for i, p in enumerate(paragraphs2) if EXISTING_LABEL in _para_text(p)), -1
    )
    if not (idx_new < idx_label):
        raise SystemExit(
            f"Verificação falhou: parágrafo do '{NEW_TOKEN}' (idx={idx_new}) "
            f"não está antes do label '{EXISTING_LABEL}' (idx={idx_label})."
        )

    print(
        f"[ok] Verificações pós-save: '{NEW_TOKEN}' presente 1×; "
        f"'{EXISTING_TOKEN}' intacto; label '{EXISTING_LABEL}' intacto."
    )
    print(f"     Célula alvo: {len(paragraphs2)} parágrafos (era {len(paragraphs)}, +1).")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
