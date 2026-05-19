from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

VERMELHO = RGBColor(0xD4, 0x00, 0x00)
PRETO = RGBColor(0x1C, 0x1C, 0x1C)


def _add_text(slide, text, x, y, w, h, size=14, color=PRETO, bold=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = str(text) if text is not None else "—"
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def build_dados_cliente(prs: Presentation, values: dict) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    _add_text(slide, "Dados do Cliente", 0.5, 0.3, 12, 0.7, size=28, color=VERMELHO, bold=True)

    fields = [
        ("Nome / Razão Social", values.get("nome_razao_social")),
        ("CNPJ / CPF", values.get("cpf_cnpj")),
        ("Contato", values.get("nome_contato")),
        ("Telefone", values.get("telefone")),
        ("E-mail", values.get("email")),
        ("Local da obra", values.get("endereco_obra")),
    ]
    for i, (label, value) in enumerate(fields):
        row_y = 1.3 + i * 0.9
        _add_text(slide, label, 0.5, row_y, 5, 0.5, size=11, color=PRETO, bold=True)
        _add_text(slide, value or "—", 6.0, row_y, 6.8, 0.5, size=13, color=PRETO)


def build_sumario(prs: Presentation, product_groups: list) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    _add_text(slide, "Sumário", 0.5, 0.3, 12, 0.7, size=28, color=VERMELHO, bold=True)

    fixed_items = [
        "2. Investimento",
        "3. Condições de Pagamento",
        "4. Prazos",
        "5. Garantia",
        "6. Responsabilidades da Contratada",
        "7. Responsabilidades do Contratante",
        "8. Considerações Gerais",
    ]

    y = 1.3
    for group in product_groups:
        _add_text(
            slide,
            f"1. Projeto — {group.sumarioText}",
            0.5, y, 12.3, 1.0,
            size=11, color=PRETO,
        )
        y += 1.0

    for item in fixed_items:
        _add_text(slide, item, 0.5, y, 12.3, 0.55, size=11, color=PRETO)
        y += 0.55


def build_investimento(prs: Presentation, group) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    title = f"Investimento — {group.productId.replace('_', ' ').title()}"
    _add_text(slide, title, 0.5, 0.3, 12, 0.7, size=28, color=VERMELHO, bold=True)

    if not group.investimentoRows:
        return

    n_rows = len(group.investimentoRows)
    table = slide.shapes.add_table(
        n_rows, 2,
        Inches(0.5), Inches(1.4),
        Inches(12.3), Inches(n_rows * 0.55),
    ).table

    table.columns[0].width = Inches(9.5)
    table.columns[1].width = Inches(2.8)

    for i, label in enumerate(group.investimentoRows):
        cell = table.cell(i, 0)
        cell.text = label
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(12)
        run.font.color.rgb = PRETO
        table.cell(i, 1).text = ""
